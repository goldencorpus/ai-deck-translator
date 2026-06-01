#!/usr/bin/env python3
"""
Simple test script for slide notes translation.

This script creates a simple PowerPoint presentation with slide notes,
simulates translation, and updates the notes.
"""
import os
from pptx import Presentation

# Create test directory if it doesn't exist
TEST_DIR = "test_output"
os.makedirs(TEST_DIR, exist_ok=True)

# File paths
TEST_PPTX = os.path.join(TEST_DIR, "notes_test.pptx")
OUTPUT_PPTX = os.path.join(TEST_DIR, "notes_test_translated.pptx")


def create_test_presentation():
    """Create a simple PowerPoint presentation with slide notes."""
    print("Creating test presentation...")
    prs = Presentation()

    # Slide 1: Title slide with notes
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Test Presentation"
    subtitle.text = "Testing Slide Notes Translation"

    # Add notes to slide 1
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the title slide."

    # Slide 2: Content slide with notes
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Content Slide"
    content.text = "This is a bullet point\nThis is another bullet point"

    # Add notes to slide 2
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the content slide."

    # Save the presentation
    prs.save(TEST_PPTX)
    print(f"Test presentation created: {TEST_PPTX}")


def extract_text(prs):
    """Extract text from a PowerPoint presentation, including slide notes."""
    print("\nExtracting text...")

    # Dictionary to store text elements
    text_elements = {}

    # List to store slide metadata
    slide_metadata = []

    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1
        print(f"Processing slide {slide_number}...")

        # Create slide metadata
        metadata = {"slide_number": slide_number, "elements": []}

        # Extract text from shapes
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = f"slide{slide_number}_shape{shape_idx}"

            # Extract text from text frames
            if hasattr(shape, "text") and shape.text.strip():
                text_elements[shape_id] = shape.text
                print(f"  Shape {shape_idx}: {shape.text}")
                metadata["elements"].append(
                    {"id": shape_id, "type": "shape", "text": shape.text}
                )

        # Extract slide notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                metadata["notes"] = notes_text.strip()
                print(f"  Notes: {notes_text.strip()}")

        # Add slide metadata to the list
        slide_metadata.append(metadata)

    return text_elements, slide_metadata


def mock_translate(text, target_language):
    """Mock translation function that adds a prefix to the text."""
    prefix_map = {
        "es": "[ES] ",
        "fr": "[FR] ",
        "de": "[DE] ",
        "ja": "[JA] ",
    }
    prefix = prefix_map.get(target_language, "[??] ")
    return prefix + text


def update_presentation(prs, translated_elements):
    """Update the presentation with translated text."""
    print("\nUpdating presentation...")

    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1
        print(f"Updating slide {slide_number}...")

        # Update text in shapes
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = f"slide{slide_number}_shape{shape_idx}"

            # Update text in text frames
            if (
                hasattr(shape, "text")
                and shape.text.strip()
                and shape_id in translated_elements
            ):
                shape.text = translated_elements[shape_id]
                print(f"  Updated shape {shape_idx}: {shape.text}")

        # Update slide notes
        notes_id = f"slide{slide_number}_notes"
        if notes_id in translated_elements and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = translated_elements[notes_id]
            print(f"  Updated notes: {slide.notes_slide.notes_text_frame.text}")


def main():
    """Main function."""
    # Create test presentation
    create_test_presentation()

    # Open the presentation
    prs = Presentation(TEST_PPTX)

    # Extract text from the presentation
    text_elements, slide_metadata = extract_text(prs)

    # Translate text elements and notes
    print("\nTranslating text elements and notes...")
    target_language = "es"  # Spanish

    # Create dictionary of translated elements
    translated_elements = {}

    # Translate text elements
    for key, value in text_elements.items():
        translated_elements[key] = mock_translate(value, target_language)
        print(f"  {key}: {translated_elements[key]}")

    # Translate notes
    for slide in slide_metadata:
        if "notes" in slide:
            slide_number = slide["slide_number"]
            notes_id = f"slide{slide_number}_notes"
            translated_elements[notes_id] = mock_translate(
                slide["notes"], target_language
            )
            print(f"  {notes_id}: {translated_elements[notes_id]}")

    # Update the presentation with translated text
    update_presentation(prs, translated_elements)

    # Save the updated presentation
    prs.save(OUTPUT_PPTX)
    print(f"\nSaved translated presentation to {OUTPUT_PPTX}")

    # Verify the translated presentation
    print("\nVerifying translated presentation...")
    prs_translated = Presentation(OUTPUT_PPTX)

    for i, slide in enumerate(prs_translated.slides):
        slide_number = i + 1
        print(f"Slide {slide_number}:")

        # Verify shapes
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                print(f"  Shape {j}: {shape.text}")

        # Verify notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            print(f"  Notes: {notes_text}")


if __name__ == "__main__":
    main()
