"""Worker preview tests — pure pptx logic (no LLM, no network)."""

import os

# preview.py imports the engine, whose config validates this at import on some paths.
os.environ.setdefault("CLAUDE_API_KEY", "test-dummy")
os.environ.setdefault("ANTHROPIC_API_KEY", "test-dummy")

import tempfile

from pptx import Presentation

from worker.preview import pick_representative_slide, isolate_slide

SAMPLE = os.path.join(os.path.dirname(__file__), "..", "test_output", "test_presentation.pptx")


def test_pick_representative_returns_valid_index():
    prs = Presentation(SAMPLE)
    n = len(prs.slides.__iter__.__self__._sldIdLst)  # slide count
    idx = pick_representative_slide(SAMPLE)
    assert isinstance(idx, int)
    assert 0 <= idx < n


def test_pick_prefers_content_over_title():
    # The sample's slide 0 is a title slide; a content slide should win.
    idx = pick_representative_slide(SAMPLE)
    prs = Presentation(SAMPLE)
    chosen = list(prs.slides)[idx]
    text = " ".join(s.text_frame.text for s in chosen.shapes if s.has_text_frame)
    assert len(text.strip()) > 0


def test_isolate_slide_yields_single_slide():
    with tempfile.TemporaryDirectory() as tmp:
        out = os.path.join(tmp, "one.pptx")
        isolate_slide(SAMPLE, 1, out)
        assert len(list(Presentation(out).slides)) == 1
