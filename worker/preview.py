"""
Watermarked preview — the free, pre-payment proof surface.

Translate ONE representative slide of the uploaded deck, render it to a PNG, stamp a
diagonal "PREVIEW" watermark, and hand it back. The buyer sees fidelity on their OWN
content before paying; the full deck stays paywalled. Cheap (~1 slide of translation).

Pipeline: pick representative slide -> isolate it into a 1-slide deck -> translate that
deck (engine) -> LibreOffice render to PNG -> Pillow watermark -> PNG bytes.
"""

from __future__ import annotations

import copy
import os
import subprocess
import tempfile

from pptx import Presentation
from pptx.util import Emu

from PIL import Image, ImageDraw, ImageFont

from ai_deck_translator.pptx.translator import translate_pptx

SOFFICE = "/usr/bin/soffice"


def pick_representative_slide(path: str) -> int:
    """First content-heavy slide (skip pure title / section dividers). Falls back to 0."""
    prs = Presentation(path)
    best_idx, best_chars = 0, -1
    for i, slide in enumerate(prs.slides):
        chars = 0
        shape_count = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                chars += len(t)
                if t:
                    shape_count += 1
        # Prefer a slide with real body text AND more than one text block (not a title slide).
        score = chars + (50 if shape_count >= 2 else 0)
        if score > best_chars:
            best_idx, best_chars = i, score
    return best_idx


def isolate_slide(src: str, keep_idx: int, dest: str) -> None:
    """Write a 1-slide deck containing only slide `keep_idx` (drop the other sldId refs)."""
    prs = Presentation(src)
    sld_id_lst = prs.slides._sldIdLst
    for i, sld_id in enumerate(list(sld_id_lst)):
        if i != keep_idx:
            sld_id_lst.remove(sld_id)
    prs.save(dest)


def render_png(pptx_path: str, out_dir: str) -> str:
    """LibreOffice headless: render the (single-slide) deck to a PNG. Returns the PNG path."""
    subprocess.run(
        [SOFFICE, "--headless", "--convert-to", "png", "--outdir", out_dir, pptx_path],
        check=True,
        capture_output=True,
        timeout=120,
        env={**os.environ, "HOME": out_dir},  # soffice needs a writable HOME profile
    )
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    png = os.path.join(out_dir, base + ".png")
    if not os.path.exists(png):
        raise RuntimeError("LibreOffice did not produce a PNG")
    return png


def watermark(png_path: str, out_path: str, text: str = "SlideVerso  ·  PREVIEW") -> None:
    """Stamp a tiled diagonal watermark so the image proves fidelity but isn't usable."""
    img = Image.open(png_path).convert("RGBA")
    w, h = img.size
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    size = max(20, w // 28)
    try:
        font = ImageFont.truetype("/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf", size)
    except OSError:
        font = ImageFont.load_default()
    # Diagonal repeating watermark across the whole slide.
    step_x, step_y = w // 2, h // 4
    stamp = Image.new("RGBA", (w, max(size * 2, 40)), (0, 0, 0, 0))
    sd = ImageDraw.Draw(stamp)
    sd.text((0, 0), text, font=font, fill=(120, 120, 120, 90))
    stamp = stamp.rotate(30, expand=True)
    for yy in range(-step_y, h + step_y, step_y):
        for xx in range(-step_x, w + step_x, step_x):
            overlay.alpha_composite(stamp, (xx, yy))
    out = Image.alpha_composite(img, overlay).convert("RGB")
    out.save(out_path, "PNG")


def make_preview(input_pptx: str, source_lang: str, target_lang: str, work: str) -> str:
    """Full pipeline → returns the path to the watermarked preview PNG."""
    idx = pick_representative_slide(input_pptx)
    one = os.path.join(work, "one.pptx")
    isolate_slide(input_pptx, idx, one)
    translated = os.path.join(work, "one_translated.pptx")
    translate_pptx(one, translated, source_lang, target_lang)
    raw_png = render_png(translated, work)
    wm = os.path.join(work, "preview.png")
    watermark(raw_png, wm)
    return wm
