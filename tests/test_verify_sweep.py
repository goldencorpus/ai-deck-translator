"""
Deterministic tests for the P2 verify pass: the no-LLM sweep detectors, the surgical patch
(critic call seam stubbed), and a full extract -> contract -> JSONL -> sweep -> patch ->
reassemble end-to-end with a deterministic fake translator. No API key required.
"""

import json
import os
import re
import shutil
import tempfile
import unittest
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches

from ai_deck_translator.pptx import verify as verify_mod
from ai_deck_translator.pptx import translator as pptx_translator
from ai_deck_translator.pptx.verify import sweep, patch as patch_translations
from ai_deck_translator.pptx.translator import translate_pptx, missing_block_ids
from ai_deck_translator.pptx.extractor import extract_text

CONTRACT = {
    "doc_context": "QBR deck",
    "register": {"style": "teineigo (です・ます)", "rules": []},
    "deixis": {},
    "glossary": {"Synergy": "シナジー"},
    "proper_nouns": {"Acme": {"canonical": "アクメ", "forbidden": ["アクミー"]}},
    "header_backbone": {},
}


class TestSweepDetectors(unittest.TestCase):
    def test_locked_term_drift(self):
        src = {"slide1_shape1": "Our Synergy strategy"}
        tgt = {
            "slide1_shape1": "我々の相乗効果戦略です"
        }  # used wrong term, not シナジー
        v = sweep(src, tgt, CONTRACT, target_language="ja")
        kinds = {x.kind for x in v}
        self.assertIn("locked_term", kinds)

    def test_locked_term_satisfied_is_clean(self):
        src = {"slide1_shape1": "Our Synergy strategy"}
        tgt = {"slide1_shape1": "我々のシナジー戦略です"}
        quality = [
            x
            for x in sweep(src, tgt, CONTRACT, target_language="ja")
            if x.kind in verify_mod._FIXABLE_KINDS
        ]
        self.assertEqual(quality, [])

    def test_forbidden_proper_noun_variant(self):
        src = {"slide1_shape1": "Acme grew"}
        tgt = {"slide1_shape1": "アクミーは成長しました"}  # forbidden variant
        kinds = {x.kind for x in sweep(src, tgt, CONTRACT, target_language="ja")}
        self.assertIn("proper_noun", kinds)

    def test_residual_source(self):
        src = {"slide1_shape1": "Quarterly revenue"}
        tgt = {"slide1_shape1": "Quarterly な収益です"}  # left English word in
        kinds = {x.kind for x in sweep(src, tgt, CONTRACT, target_language="ja")}
        self.assertIn("residual_source", kinds)

    def test_keigo_role_aware(self):
        src = {"slide1_shape0": "Big Title", "slide1_shape1": "Body claim"}
        tgt = {
            "slide1_shape0": "事業である",  # plain ending — but this is a TITLE -> exempt
            "slide1_shape1": "これは事実である",  # plain ending in body -> flagged
        }
        meta = {
            "slide1_shape0": {"role": "title"},
            "slide1_shape1": {"role": "body_bullet"},
        }
        v = sweep(src, tgt, CONTRACT, blocks_meta=meta, target_language="ja")
        keigo = [x for x in v if x.kind == "keigo"]
        flagged = {x.id for x in keigo}
        self.assertIn("slide1_shape1", flagged)
        self.assertNotIn("slide1_shape0", flagged)  # title exempt from です・ます

    def test_completeness_and_orphans(self):
        src = {"a": "x", "b": "y"}
        tgt = {"a": "翻訳", "c": "orphan"}  # b missing, c orphan
        v = sweep(src, tgt, {}, target_language="ja")
        kinds = {(x.id, x.kind) for x in v}
        self.assertIn(("b", "missing"), kinds)
        self.assertIn(("c", "orphan"), kinds)

    def test_empty_contract_skips_quality_checks(self):
        src = {"slide1_shape1": "Synergy"}
        tgt = {"slide1_shape1": "Synergy"}  # would be residual, but no contract
        quality = [
            x
            for x in sweep(src, tgt, {}, target_language="ja")
            if x.kind in verify_mod._FIXABLE_KINDS
        ]
        self.assertEqual(quality, [])


class TestPatch(unittest.TestCase):
    def test_patch_fixes_flagged_block(self):
        src = {"slide1_shape1": "Our Synergy strategy"}
        tgt = {"slide1_shape1": "我々の相乗効果戦略です"}
        violations = sweep(src, tgt, CONTRACT, target_language="ja")

        def fake_complete(system, user, **kwargs):
            self.assertIn(
                "シナジー", system
            )  # contract rode into the critic system prompt
            return json.dumps(
                {"id": "slide1_shape1", "fix": "我々のシナジー戦略です"},
                ensure_ascii=False,
            )

        with patch.object(verify_mod, "_complete", side_effect=fake_complete):
            result = patch_translations(violations, src, tgt, CONTRACT, "en", "ja")
        self.assertEqual(result["slide1_shape1"], "我々のシナジー戦略です")

    def test_patch_degrades_on_critic_failure(self):
        src = {"slide1_shape1": "Our Synergy strategy"}
        tgt = {"slide1_shape1": "我々の相乗効果戦略です"}
        violations = sweep(src, tgt, CONTRACT, target_language="ja")
        with patch.object(verify_mod, "_complete", side_effect=RuntimeError("boom")):
            result = patch_translations(violations, src, tgt, CONTRACT, "en", "ja")
        self.assertEqual(result, tgt)  # unchanged, never blocks

    def test_patch_noop_when_no_fixable_violations(self):
        src = {"a": "x"}
        tgt = {"a": "翻訳"}
        with patch.object(verify_mod, "_complete") as m:
            result = patch_translations([], src, tgt, CONTRACT, "en", "ja")
        m.assert_not_called()
        self.assertEqual(result, tgt)


class TestVerifyEndToEnd(unittest.TestCase):
    """extract -> contract -> fake JSONL batch (with a seeded violation) -> sweep -> patch."""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        self.input_path = os.path.join(self.tmp, "deck.pptx")
        self.output_path = os.path.join(self.tmp, "deck_ja.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Strategy Overview"
        slide.placeholders[1].text = "Our Synergy drives growth\nMarket share expanded"
        slide.notes_slide.notes_text_frame.text = "Mention the Synergy story"
        prs.save(self.input_path)

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_seeded_violation_is_patched_and_deck_is_complete(self):
        source, _ = extract_text(self.input_path)

        def fake_batch(batch, *a, **k):
            # Deterministic JA-ish output that mistranslates "Synergy" as the wrong term
            # (相乗効果) instead of the locked シナジー, and keeps no Latin elsewhere.
            out = {}
            for bid, v in batch.items():
                if "Synergy" in v:
                    out[bid] = "相乗効果が成長を促進します"
                else:
                    out[bid] = "翻訳されたテキスト"
            return out

        def fake_patch_complete(system, user, **kwargs):
            ids = list(dict.fromkeys(re.findall(r"slide\d+_\w+", user)))
            return "\n".join(
                json.dumps(
                    {"id": i, "fix": "シナジーが成長を促進します"}, ensure_ascii=False
                )
                for i in ids
            )

        with patch.object(
            pptx_translator, "translate_batch", side_effect=fake_batch
        ), patch.object(
            pptx_translator, "build_contract", return_value=CONTRACT
        ), patch.object(
            verify_mod, "_complete", side_effect=fake_patch_complete
        ):
            result = translate_pptx(self.input_path, self.output_path, "en", "ja")

        self.assertTrue(result)
        out_dict, _ = extract_text(self.output_path)
        # 100% coverage preserved.
        self.assertEqual(missing_block_ids(source, out_dict), [])
        # The seeded locked-term violation was surgically corrected to the locked form.
        synergy_blocks = [bid for bid, src in source.items() if "Synergy" in src]
        self.assertTrue(synergy_blocks)
        for bid in synergy_blocks:
            self.assertIn("シナジー", out_dict[bid])
            self.assertNotIn("相乗効果", out_dict[bid])


if __name__ == "__main__":
    unittest.main()
