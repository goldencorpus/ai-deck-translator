"""
Completeness tests for the PPTX translation pipeline.

These verify the core revenue-blocking guarantee (business-plan §8): a deck either
translates 100% of its text blocks or the job fails loudly — it never writes a
partially-translated deck. They run a real multi-slide deck (with a table and speaker
notes) through the full extract -> translate -> update -> re-extract pipeline with a
deterministic fake translator, so they need no API key and incur no cost.
"""

import os
import shutil
import tempfile
import unittest
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ai_deck_translator.pptx import translator as pptx_translator
from ai_deck_translator.pptx.translator import translate_pptx, missing_block_ids
from ai_deck_translator.pptx.extractor import extract_text
from ai_deck_translator.utils.exceptions import IncompleteTranslationError


def build_sample_deck(path):
    """Create a small multi-slide deck with a title, body, table, and speaker notes."""
    prs = Presentation()

    # Slide 1: title + multi-line body + speaker notes
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Quarterly Business Review"
    slide1.placeholders[1].text = "Revenue grew 40 percent\nCustomers exceeded 300"
    slide1.notes_slide.notes_text_frame.text = "Emphasize the growth story here."

    # Slide 2: title-only layout + a 2x2 table
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Key Metrics"
    table = slide2.shapes.add_table(
        2, 2, Inches(1), Inches(2), Inches(6), Inches(2)
    ).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Customer Satisfaction"
    table.cell(1, 1).text = "95 percent"

    prs.save(path)


def fake_translate_batch(batch, *args, **kwargs):
    """Deterministic stand-in for the Anthropic call: prefixes every value."""
    return {key: f"[ja] {value}" for key, value in batch.items()}


class TestPptxCompleteness(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        self.input_path = os.path.join(self.tmp, "deck.pptx")
        self.output_path = os.path.join(self.tmp, "deck_ja.pptx")
        build_sample_deck(self.input_path)

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_full_coverage_end_to_end(self):
        """100% of source text blocks must appear translated in the output deck."""
        source, metadata = extract_text(self.input_path)

        # Sanity: the fixture really does contain a table and speaker notes.
        self.assertTrue(any("_table_" in bid for bid in source), "no table block")
        self.assertTrue(any(bid.endswith("_notes") for bid in source), "no notes block")
        self.assertGreaterEqual(len(source), 6)

        with patch.object(
            pptx_translator, "translate_batch", side_effect=fake_translate_batch
        ):
            result = translate_pptx(self.input_path, self.output_path, "en", "ja")

        self.assertTrue(result)
        self.assertTrue(os.path.exists(self.output_path))

        out_dict, _ = extract_text(self.output_path)

        # Every single source block must be present and translated (non-blank, changed).
        for bid, original in source.items():
            self.assertIn(bid, out_dict, f"block {bid} dropped from output")
            self.assertTrue(out_dict[bid].strip(), f"block {bid} is blank in output")
            self.assertNotEqual(
                out_dict[bid].strip(),
                original.strip(),
                f"block {bid} was not translated",
            )

        # Authoritative coverage assertion: zero missing blocks.
        self.assertEqual(
            missing_block_ids(source, out_dict),
            [],
            "output deck has untranslated blocks",
        )

    def test_raises_and_writes_nothing_when_incomplete(self):
        """If any block is left untranslated, fail loudly and write no output."""
        source, _ = extract_text(self.input_path)

        # Drop the slide-1 title (position is unique, so the validator's positional
        # fallback cannot mask the gap) to simulate a genuine skip.
        dropped = "slide1_shape0"
        self.assertIn(dropped, source)
        partial = {k: f"[ja] {v}" for k, v in source.items() if k != dropped}

        with patch.object(
            pptx_translator, "translate_text", return_value=partial
        ), patch.object(pptx_translator, "update_slides") as mock_update:
            with self.assertRaises(IncompleteTranslationError) as ctx:
                translate_pptx(self.input_path, self.output_path, "en", "ja")

        self.assertIn(dropped, ctx.exception.missing_ids)
        self.assertEqual(ctx.exception.total, len(source))
        mock_update.assert_not_called()
        self.assertFalse(
            os.path.exists(self.output_path),
            "a partial deck must never be written",
        )


class TestPptxFormattingPreserved(unittest.TestCase):
    """Translating must keep run formatting and never leak source-language fragments."""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        self.input_path = os.path.join(self.tmp, "fmt.pptx")
        self.output_path = os.path.join(self.tmp, "fmt_ja.pptx")

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Mixed-run paragraph: "Revenue " + bold green "grew strongly" + " this year"
        box = slide.shapes.add_textbox(Pt(50), Pt(80), Pt(500), Pt(80)).text_frame
        para = box.paragraphs[0]
        for text, bold, color in [
            ("Revenue ", None, None),
            ("grew strongly", True, RGBColor(0x00, 0x70, 0x00)),
            (" this year", None, None),
        ]:
            run = para.add_run()
            run.text = text
            run.font.size = Pt(24)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color

        # Styled table header cell.
        table = slide.shapes.add_table(
            2, 1, Pt(50), Pt(200), Pt(300), Pt(80)
        ).table
        header = table.cell(0, 0)
        header.text = "Metric"
        hrun = header.text_frame.paragraphs[0].runs[0]
        hrun.font.bold = True
        hrun.font.size = Pt(16)
        hrun.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        table.cell(1, 0).text = "Growth"

        prs.save(self.input_path)

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_no_leftover_runs_and_table_format_preserved(self):
        def fake(batch, *args, **kwargs):
            return {key: "JP_" + value for key, value in batch.items()}

        with patch.object(pptx_translator, "translate_batch", side_effect=fake):
            translate_pptx(self.input_path, self.output_path, "en", "ja")

        prs = Presentation(self.output_path)
        slide = prs.slides[0]

        # The mixed-run paragraph must collapse to a single run with no English left.
        textbox = next(
            s
            for s in slide.shapes
            if s.has_text_frame and "Revenue" in s.text_frame.text
        )
        runs = textbox.text_frame.paragraphs[0].runs
        self.assertEqual(len(runs), 1, "leftover runs would leak source-language text")
        # Exact text (no duplication): the old bug left the original runs in place,
        # producing "JP_Revenue grew strongly this yeargrew strongly this year".
        self.assertEqual(
            textbox.text_frame.text, "JP_Revenue grew strongly this year"
        )

        # The styled table header must keep its bold / size / colour.
        table = next(s for s in slide.shapes if s.has_table).table
        hrun = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        self.assertEqual(hrun.text, "JP_Metric")
        self.assertTrue(hrun.font.bold)
        self.assertEqual(hrun.font.size, Pt(16))
        self.assertEqual(hrun.font.color.rgb, RGBColor(0xFF, 0xFF, 0xFF))


class TestMultiParagraphStructure(unittest.TestCase):
    """A multi-bullet text box must stay multi-bullet — not collapse into one line."""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        self.input_path = os.path.join(self.tmp, "bullets.pptx")
        self.output_path = os.path.join(self.tmp, "bullets_ja.pptx")

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Pt(50), Pt(100), Pt(600), Pt(300)).text_frame
        box.paragraphs[0].text = "First point about speed"
        for line in ("Second point about cost", "Third point about trust"):
            box.add_paragraph().text = line
        prs.save(self.input_path)

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_bullets_not_collapsed(self):
        # Fake translator drops newlines if ever handed a multi-line block — proving the
        # structure is preserved by per-paragraph extraction, not by echoing newlines.
        def fake(batch, *args, **kwargs):
            return {key: "[ja] " + value.replace("\n", " ") for key, value in batch.items()}

        source, _ = extract_text(self.input_path)
        # Each bullet should be its own extracted block.
        self.assertEqual(sum(1 for k in source if "_p" in k), 3)

        with patch.object(pptx_translator, "translate_batch", side_effect=fake):
            translate_pptx(self.input_path, self.output_path, "en", "ja")

        prs = Presentation(self.output_path)
        box = next(
            s
            for s in prs.slides[0].shapes
            if s.has_text_frame and "[ja]" in s.text_frame.text
        )
        non_empty = [p for p in box.text_frame.paragraphs if p.text.strip()]
        self.assertEqual(
            len(non_empty), 3, "the three bullets must remain three separate paragraphs"
        )


if __name__ == "__main__":
    unittest.main()
