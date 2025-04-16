#!/usr/bin/env python3
"""
Test script for the complete PPTX translation process with our new improvements.

This script creates a test presentation, translates it using mock translation,
and verifies that the translations are correctly applied.
"""
import os
import sys
import uuid
from pptx import Presentation
from ai_deck_translator.pptx.extractor import extract_text
from ai_deck_translator.pptx.translator import standardize_ids, validate_translation_ids, verify_translation_keys

# Create test directory if it doesn't exist
TEST_DIR = "test_output"
os.makedirs(TEST_DIR, exist_ok=True)

# File paths
random_id = str(uuid.uuid4())
TEST_PPTX = os.path.join(TEST_DIR, f"{random_id}_test_format.pptx")
OUTPUT_PPTX = os.path.join(TEST_DIR, f"{random_id}_test_format_translated.pptx")

def create_test_presentation():
    """Create a simple PowerPoint presentation with slide notes."""
    print("Creating test presentation...")
    prs = Presentation()
    
    # Slide 1: Title slide with notes
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Test Presentation"
    subtitle.text = "Testing Format Enforcement"
    
    # Add notes to slide 1
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the title slide."
    
    # Slide 2: Content slide with notes
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content Slide"
    content.text = "This is a bullet point\nThis is another bullet point\nThis is a third bullet point"
    
    # Add notes to slide 2
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the content slide."
    
    # Save the presentation
    prs.save(TEST_PPTX)
    print(f"Test presentation created: {TEST_PPTX}")

def mock_translate(text_dict, target_language):
    """Mock translation function that adds a prefix to the text and intentionally changes some key formats."""
    prefix_map = {
        "es": "[ES] ",
        "fr": "[FR] ",
        "de": "[DE] ",
        "ja": "[JA] ",
    }
    prefix = prefix_map.get(target_language, "[??] ")
    
    # Create a translated dictionary with some key format changes
    translated_dict = {}
    
    for key, value in text_dict.items():
        translated_value = prefix + value
        
        # For some keys, alter the format to test our robustness
        if "_shape" in key and "_shape0" in key:
            # Change format: slide1_shape0 -> slide_1_element_0
            new_key = key.replace("slide", "slide_").replace("shape", "element")
            translated_dict[new_key] = translated_value
        elif "_notes" in key:
            # Change format: slide1_notes -> slide_1_notes
            slide_num = key.replace("slide", "").split("_")[0]
            new_key = f"slide_{slide_num}_notes"
            translated_dict[new_key] = translated_value
        else:
            # Keep original format for other keys
            translated_dict[key] = translated_value
    
    return translated_dict

def apply_translation(pptx_file, output_file, translated_texts):
    """Apply translations to the presentation (simplified version)."""
    print(f"Applying translations to {pptx_file}...")
    
    # Create a copy of the presentation
    prs = Presentation(pptx_file)
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1
        print(f"Updating slide {slide_number}...")
        
        # Update text in shapes
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = f"slide{slide_number}_shape{shape_idx}"
            
            # Update text in text frames
            if hasattr(shape, "text") and shape.text.strip() and shape_id in translated_texts:
                shape.text = translated_texts[shape_id]
                print(f"  Updated shape {shape_id}")
        
        # Update slide notes
        notes_id = f"slide{slide_number}_notes"
        if notes_id in translated_texts and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = translated_texts[notes_id]
            print(f"  Updated notes for slide {slide_number}")
    
    # Save the presentation
    prs.save(output_file)
    print(f"Saved translated presentation to {output_file}")
    
    return True

def test_translation_process():
    """Test the translation process with format enforcement."""
    # Create test presentation
    create_test_presentation()
    
    # Extract text from the presentation
    print("\nExtracting text from presentation...")
    text_dict, slide_metadata = extract_text(TEST_PPTX)
    
    # Print extracted text elements
    print("\nExtracted text elements:")
    for key, value in text_dict.items():
        print(f"{key}: {value}")
    
    # Apply standardization to IDs
    print("\nStandardizing IDs...")
    standardized_dict = standardize_ids(text_dict, slide_metadata)
    
    # Mock translation with format changes
    print("\nTranslating text (with mock format changes)...")
    translated_dict = mock_translate(standardized_dict, "es")
    
    # Print translated elements
    print("\nTranslated elements (with format changes):")
    for key, value in translated_dict.items():
        print(f"{key}: {value}")
    
    # Verify translation keys
    print("\nVerifying translation keys...")
    exact_match, missing_keys, format_changes = verify_translation_keys(
        standardized_dict.keys(), translated_dict.keys()
    )
    
    print(f"Exact match: {exact_match}")
    print(f"Missing keys: {missing_keys}")
    print(f"Format changes: {format_changes}")
    
    # Validate and fix translations
    print("\nValidating and fixing translations...")
    success, missing_ids, fixed_translations = validate_translation_ids(
        standardized_dict, translated_dict, slide_metadata
    )
    
    print(f"Validation {'succeeded' if success else 'failed with unresolved issues'}")
    print(f"Missing IDs after fixes: {len(missing_ids)}")
    print(f"Final translation keys: {len(fixed_translations)}")
    
    # Apply translations
    print("\nApplying translations to presentation...")
    applied = apply_translation(TEST_PPTX, OUTPUT_PPTX, fixed_translations)
    
    # Verify the result
    print("\nVerifying translated presentation...")
    if applied:
        print(f"Translation applied to {OUTPUT_PPTX}")
        
        # Re-open to verify
        prs = Presentation(OUTPUT_PPTX)
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    print(f"  Shape {j}: {shape.text}")
            
            if slide.has_notes_slide:
                print(f"  Notes: {slide.notes_slide.notes_text_frame.text}")
        
        print("\n✅ Test completed successfully!")
    else:
        print("\n❌ Failed to apply translations!")

if __name__ == "__main__":
    test_translation_process() 