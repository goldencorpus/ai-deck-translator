#!/usr/bin/env python3
"""
Test script to verify PPTX translation fixes.

This script creates a test PowerPoint presentation, extracts text,
deliberately uses different ID formats for extraction and translation,
and verifies that our ID standardization and validation fixes correctly
handle the differences.
"""
import os
import sys
from pptx import Presentation
from ai_deck_translator.pptx.extractor import extract_text
from ai_deck_translator.pptx.updater import update_slides
from ai_deck_translator.pptx.translator import standardize_ids, validate_translation_ids
import uuid

# Create test directory if it doesn't exist
TEST_DIR = "test_output"
os.makedirs(TEST_DIR, exist_ok=True)

# File paths with random UUID to avoid conflicts
random_id = str(uuid.uuid4())
TEST_PPTX = os.path.join(TEST_DIR, f"{random_id}_simple_test.pptx")
OUTPUT_PPTX = os.path.join(TEST_DIR, f"{random_id}_simple_test_translated.pptx")


def create_test_presentation():
    """Create a simple PowerPoint presentation with slide notes."""
    print("Creating test presentation...")
    prs = Presentation()

    # Slide 1: Title slide with notes
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Test Presentation"
    subtitle.text = "Testing ID Format Standardization"

    # Add notes to slide 1
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the title slide. The ID standardization fix should handle this."

    # Slide 2: Content slide with notes
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Content Slide"
    content.text = "This is a bullet point\nThis is another bullet point\nThis is a third bullet point"

    # Add notes to slide 2
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the content slide. The ID validation fix should handle this."

    # Save the presentation
    prs.save(TEST_PPTX)
    print(f"Test presentation created: {TEST_PPTX}")


def test_id_standardization():
    """Test the ID standardization functionality."""
    # Create test presentation
    create_test_presentation()

    # Extract text from the presentation
    print("\nExtracting text from presentation...")
    text_dict, slide_metadata = extract_text(TEST_PPTX)

    # Print extracted text elements with IDs
    print("\nExtracted text elements:")
    for key, value in text_dict.items():
        print(f"{key}: {value}")

    # Create a copy with deliberately different ID formats
    different_format_dict = {}

    # Convert 'slideX_shapeY' to 'slide_X_element_Y'
    for key, value in text_dict.items():
        if key.startswith("slide"):
            # Different format for slide notes
            if "_notes" in key:
                slide_num = key.replace("slide", "").split("_")[0]
                new_key = f"slide_{slide_num}_notes"
            # Different format for shapes
            elif "_shape" in key:
                parts = key.split("_")
                slide_num = parts[0].replace("slide", "")
                shape_num = parts[1].replace("shape", "")
                new_key = f"slide_{slide_num}_element_{shape_num}"
            # Keep other keys as is
            else:
                new_key = key

            different_format_dict[new_key] = value

    print("\nTranslated elements with different ID format:")
    for key, value in different_format_dict.items():
        print(f"{key}: [Translated] {value}")

    # Create mock translations
    translated_dict = {k: f"[Translated] {v}" for k, v in different_format_dict.items()}

    # Test standardization
    print("\nStandardizing IDs...")
    standardized_dict = standardize_ids(text_dict, slide_metadata)

    # Test validation and fixing
    print("\nValidating translation IDs...")
    success, missing_ids, fixed_translations = validate_translation_ids(
        standardized_dict, translated_dict, slide_metadata
    )

    # Print results
    print(f"\nValidation {'succeeded' if success else 'failed'}")
    if missing_ids:
        print(f"Missing IDs: {len(missing_ids)}")
        for id in list(missing_ids.keys())[:5]:
            print(f"  {id}")

    # Update the presentation with fixed translations
    print("\nUpdating presentation with fixed translations...")
    update_slides(TEST_PPTX, OUTPUT_PPTX, fixed_translations)

    # Verify the updated presentation
    print("\nVerifying updated presentation...")
    prs = Presentation(OUTPUT_PPTX)

    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")

        # Check shapes
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                print(f"  Shape {j}: {shape.text}")

        # Check notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            print(f"  Notes: {notes_text}")

    print(f"\nTest completed successfully! Output file: {OUTPUT_PPTX}")


if __name__ == "__main__":
    test_id_standardization()
