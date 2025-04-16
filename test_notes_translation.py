#!/usr/bin/env python3
"""
Test script for slide notes translation feature.

This script creates a simple PowerPoint presentation with slide notes,
extracts text from the presentation, translates it, and updates the
presentation with the translated text.
"""
import os
import sys
from pptx import Presentation
from ai_deck_translator.pptx.extractor import extract_text
from ai_deck_translator.pptx.updater import update_slides

# Create test directory if it doesn't exist
TEST_DIR = "test_output"
os.makedirs(TEST_DIR, exist_ok=True)

# File paths
TEST_PPTX = os.path.join(TEST_DIR, "test_presentation.pptx")
OUTPUT_PPTX = os.path.join(TEST_DIR, "test_presentation_translated.pptx")

def create_test_presentation():
    """Create a simple PowerPoint presentation with slide notes."""
    print("Creating test presentation...")
    prs = Presentation()
    
    # Slide 1: Title slide with notes
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Test Presentation"
    subtitle.text = "Testing Slide Notes Translation"
    
    # Add notes to slide 1
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the title slide. They should be translated."
    
    # Slide 2: Content slide with notes
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content Slide"
    content.text = "This is a bullet point\nThis is another bullet point\nThis is a third bullet point"
    
    # Add notes to slide 2
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the content slide. They contain important information for the presenter."
    
    # Slide 3: Content slide with notes
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Another Content Slide"
    content.text = "This slide has different content\nWith multiple bullet points\nAnd some more text"
    
    # Add notes to slide 3
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are more detailed speaker notes. They explain what the presenter should say during this slide."
    
    # Save the presentation
    prs.save(TEST_PPTX)
    print(f"Test presentation created: {TEST_PPTX}")

def mock_translate(text, target_language):
    """Mock translation function that adds a prefix to the text."""
    prefix_map = {
        "es": "[ES] ",
        "fr": "[FR] ",
        "de": "[DE] ",
        "ja": "[JA] ",
    }
    prefix = prefix_map.get(target_language, "[??] ")
    return prefix + text

def test_notes_translation():
    """Test the slide notes translation feature."""
    # Create test presentation
    create_test_presentation()
    
    # Extract text from the presentation
    print("\nExtracting text from presentation...")
    text_elements, slide_metadata = extract_text(TEST_PPTX)
    
    # Print extracted text elements
    print("\nExtracted text elements:")
    for key, value in text_elements.items():
        print(f"{key}: {value}")
    
    # Print slide metadata (including notes)
    print("\nSlide metadata (including notes):")
    for slide in slide_metadata:
        print(f"Slide {slide['slide_number']}:")
        if 'notes' in slide and slide['notes']:
            print(f"  Notes: {slide['notes']}")
    
    # Translate text elements and notes
    print("\nTranslating text elements and notes...")
    target_language = "es"  # Spanish
    
    # Create dictionary of translated elements
    translated_elements = {}
    
    # Translate text elements
    for key, value in text_elements.items():
        translated_elements[key] = mock_translate(value, target_language)
    
    # Translate notes
    for slide in slide_metadata:
        if 'notes' in slide and slide['notes']:
            slide_number = slide['slide_number']
            notes_id = f"slide{slide_number}_notes"
            translated_elements[notes_id] = mock_translate(slide['notes'], target_language)
    
    # Print translated elements
    print("\nTranslated elements:")
    for key, value in translated_elements.items():
        print(f"{key}: {value}")
    
    # Update the presentation with translated text
    print("\nUpdating presentation with translated text...")
    success = update_slides(TEST_PPTX, OUTPUT_PPTX, translated_elements)
    
    if success:
        print(f"\nTranslation successful! Output file: {OUTPUT_PPTX}")
    else:
        print("\nTranslation failed!")

if __name__ == "__main__":
    test_notes_translation() 