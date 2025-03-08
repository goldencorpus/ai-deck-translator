#!/usr/bin/env python3
from pptx import Presentation
import re
import sys

def has_japanese(text):
    """Check if text contains Japanese characters"""
    japanese_range = re.compile(r'[\u3000-\u303f\u3040-\u309f\u30a0-\u30ff\uff00-\uff9f\u4e00-\u9faf]')
    return bool(japanese_range.search(text))

def has_english(text):
    """Check if text contains English characters"""
    english_range = re.compile(r'[a-zA-Z]')
    return bool(english_range.search(text))

def analyze_pptx(file_path):
    """Analyze a PowerPoint file's text content and formatting"""
    prs = Presentation(file_path)
    slide_count = len(prs.slides)
    
    slides_data = []
    total_text_elements = 0
    japanese_elements = 0
    english_elements = 0
    mixed_elements = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "slide_num": slide_idx + 1,
            "shapes": [],
            "tables": []
        }
        
        # Process shapes
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text_frame") and shape.text.strip():
                total_text_elements += 1
                text = shape.text.strip()
                
                is_japanese = has_japanese(text)
                is_english = has_english(text)
                
                if is_japanese and is_english:
                    text_type = "mixed"
                    mixed_elements += 1
                elif is_japanese:
                    text_type = "japanese"
                    japanese_elements += 1
                elif is_english:
                    text_type = "english"
                    english_elements += 1
                else:
                    text_type = "other"
                
                # Get font sizes from runs
                font_sizes = []
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if hasattr(run, "font") and hasattr(run.font, "size") and run.font.size:
                                font_sizes.append(run.font.size)
                
                shape_info = {
                    "id": f"slide_{slide_idx+1}_shape_{shape_idx}",
                    "text": text[:50] + ("..." if len(text) > 50 else ""),
                    "text_type": text_type,
                    "font_sizes": font_sizes
                }
                
                slide_data["shapes"].append(shape_info)
            
            # Check tables
            if hasattr(shape, "has_table") and shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            total_text_elements += 1
                            text = cell.text.strip()
                            
                            is_japanese = has_japanese(text)
                            is_english = has_english(text)
                            
                            if is_japanese and is_english:
                                text_type = "mixed"
                                mixed_elements += 1
                            elif is_japanese:
                                text_type = "japanese"
                                japanese_elements += 1
                            elif is_english:
                                text_type = "english"
                                english_elements += 1
                            else:
                                text_type = "other"
                            
                            # Get font sizes from cell text
                            font_sizes = []
                            if hasattr(cell, "text_frame"):
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if hasattr(run, "font") and hasattr(run.font, "size") and run.font.size:
                                            font_sizes.append(run.font.size)
                            
                            cell_info = {
                                "id": f"slide_{slide_idx+1}_table_{shape_idx}_r{row_idx}_c{col_idx}",
                                "text": text[:50] + ("..." if len(text) > 50 else ""),
                                "text_type": text_type,
                                "font_sizes": font_sizes
                            }
                            
                            slide_data["tables"].append(cell_info)
        
        slides_data.append(slide_data)
    
    return {
        "slide_count": slide_count,
        "total_text_elements": total_text_elements,
        "japanese_elements": japanese_elements,
        "english_elements": english_elements,
        "mixed_elements": mixed_elements,
        "slides": slides_data
    }

def compare_pptx_files(before_path, after_path):
    """Compare before and after PowerPoint files"""
    before_data = analyze_pptx(before_path)
    after_data = analyze_pptx(after_path)
    
    print(f"=== Before Translation ===")
    print(f"Slide count: {before_data['slide_count']}")
    print(f"Total text elements: {before_data['total_text_elements']}")
    print(f"English elements: {before_data['english_elements']}")
    print(f"Japanese elements: {before_data['japanese_elements']}")
    print(f"Mixed elements: {before_data['mixed_elements']}")
    print()
    
    print(f"=== After Translation ===")
    print(f"Slide count: {after_data['slide_count']}")
    print(f"Total text elements: {after_data['total_text_elements']}")
    print(f"English elements: {after_data['english_elements']}")
    print(f"Japanese elements: {after_data['japanese_elements']}")
    print(f"Mixed elements: {after_data['mixed_elements']}")
    print()
    
    # Calculate translation rate
    if before_data['english_elements'] > 0:
        translation_rate = (after_data['japanese_elements'] + after_data['mixed_elements']) / before_data['english_elements'] * 100
        print(f"Translation rate: {translation_rate:.1f}%")
    
    # Check for font size issues
    font_issues = []
    for b_slide, a_slide in zip(before_data["slides"], after_data["slides"]):
        for b_shape in b_slide["shapes"]:
            for a_shape in a_slide["shapes"]:
                if b_shape["id"] == a_shape["id"] and b_shape["font_sizes"] and a_shape["font_sizes"]:
                    if b_shape["font_sizes"] != a_shape["font_sizes"]:
                        font_issues.append({
                            "id": b_shape["id"],
                            "before_text": b_shape["text"],
                            "after_text": a_shape["text"],
                            "before_sizes": b_shape["font_sizes"],
                            "after_sizes": a_shape["font_sizes"]
                        })
    
    print(f"\n=== Font Size Issues ===")
    print(f"Found {len(font_issues)} text elements with changed font sizes")
    
    if font_issues:
        for i, issue in enumerate(font_issues[:5]):  # Show first 5 issues
            print(f"\nIssue {i+1}:")
            print(f"  Element: {issue['id']}")
            print(f"  Before text: {issue['before_text']}")
            print(f"  After text: {issue['after_text']}")
            print(f"  Before sizes: {issue['before_sizes']}")
            print(f"  After sizes: {issue['after_sizes']}")
        
        if len(font_issues) > 5:
            print(f"\n... and {len(font_issues) - 5} more issues")
    
    # Find untranslated elements
    untranslated = []
    for slide_idx, slide in enumerate(after_data["slides"]):
        for shape in slide["shapes"]:
            if shape["text_type"] == "english":
                untranslated.append({
                    "id": shape["id"],
                    "text": shape["text"],
                    "slide": slide_idx + 1
                })
        
        for cell in slide["tables"]:
            if cell["text_type"] == "english":
                untranslated.append({
                    "id": cell["id"],
                    "text": cell["text"],
                    "slide": slide_idx + 1
                })
    
    print(f"\n=== Untranslated Elements ===")
    print(f"Found {len(untranslated)} untranslated text elements")
    
    if untranslated:
        for i, element in enumerate(untranslated[:10]):  # Show first 10 untranslated
            print(f"\nUntranslated {i+1}:")
            print(f"  Element: {element['id']}")
            print(f"  Slide: {element['slide']}")
            print(f"  Text: {element['text']}")
        
        if len(untranslated) > 10:
            print(f"\n... and {len(untranslated) - 10} more untranslated elements")
    
    return {
        "before": before_data,
        "after": after_data,
        "font_issues": font_issues,
        "untranslated": untranslated
    }

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python compare_pptx.py <before_file.pptx> <after_file.pptx>")
        sys.exit(1)
    
    before_path = sys.argv[1]
    after_path = sys.argv[2]
    
    compare_pptx_files(before_path, after_path)