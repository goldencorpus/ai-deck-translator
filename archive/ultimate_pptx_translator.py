#!/usr/bin/env python3
import os
import json
import anthropic
import sys
import re
import time
import argparse
from datetime import datetime
from tqdm import tqdm
from dotenv import load_dotenv
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation

# Load environment variables from .env file
load_dotenv()

# XML namespaces used in PPTX files
namespaces = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
    'a14': 'http://schemas.microsoft.com/office/drawing/2010/main',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
    'sl': 'http://schemas.openxmlformats.org/officeDocument/2006/slideLayout',
    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
}

def extract_text_from_element(element):
    """Extract text from an XML element and its children."""
    text = ""
    
    # Extract text from a:t elements (text runs)
    for t in element.findall('.//a:t', namespaces):
        if t.text:
            text += t.text + " "
    
    # Extract text from r:t elements (some older PPT formats)
    for t in element.findall('.//r:t', namespaces):
        if t.text:
            text += t.text + " "
            
    # Extract text from mc:t elements (compat mode text)
    for t in element.findall('.//mc:t', namespaces):
        if t.text:
            text += t.text + " "
            
    # Also check direct text content of element and children
    if element.text and element.text.strip():
        text += element.text.strip() + " "
        
    for child in element:
        if child.text and child.text.strip():
            text += child.text.strip() + " "
        if child.tail and child.tail.strip():
            text += child.tail.strip() + " "
    
    return text.strip()

def extract_from_smartart(pptx_file, rels_path, rel_id):
    """Extract text from SmartArt diagrams using direct XML processing."""
    with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
        # Find the relationship target for the SmartArt
        rels_xml = zip_ref.read(rels_path)
        rels_root = ET.fromstring(rels_xml)
        
        # Find the target for this relationship ID
        target = None
        for rel in rels_root.findall('.//Relationship', namespaces):
            if rel.get('Id') == rel_id and rel.get('Type').endswith('diagramData'):
                target = rel.get('Target')
                break
        
        if not target:
            return ""
        
        # Convert the target path to the correct format
        if target.startswith('../'):
            target = target.replace('../', '')
        else:
            slide_dir = os.path.dirname(rels_path)
            target = os.path.join(os.path.dirname(slide_dir), target)
        
        # Read the diagram data
        try:
            diagram_xml = zip_ref.read(target)
            diagram_root = ET.fromstring(diagram_xml)
            
            # Extract text more comprehensively from the diagram
            all_text = []
            
            # Standard text elements
            for t_element in diagram_root.findall('.//a:t', namespaces):
                if t_element.text:
                    all_text.append(t_element.text)
            
            # Text in data model
            for text_elem in diagram_root.findall('.//dgm:t', {'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'}):
                if text_elem.text:
                    all_text.append(text_elem.text)
                    
            # Text in properties
            for prop in diagram_root.findall('.//dgm:p', {'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'}):
                val = prop.get('val')
                if val and not any(val in t for t in all_text):
                    all_text.append(val)
                    
            # Get the layout information - it might have more text
            data_model_target = None
            for data_rel in diagram_root.findall('.//{http://schemas.openxmlformats.org/officeDocument/2006/relationships}dm'):
                data_model_id = data_rel.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}r')
                if data_model_id:
                    # Find the data model relationship
                    dm_path = os.path.join('ppt/diagrams/_rels', os.path.basename(target) + '.rels')
                    if dm_path in zip_ref.namelist():
                        dm_rels = ET.fromstring(zip_ref.read(dm_path))
                        for rel in dm_rels.findall('.//Relationship'):
                            if rel.get('Id') == data_model_id:
                                data_model_target = rel.get('Target')
                                if data_model_target.startswith('../'):
                                    data_model_target = data_model_target.replace('../', '')
                                else:
                                    data_model_target = os.path.join('ppt/diagrams', data_model_target)
                                
                                # Extract text from the data model
                                try:
                                    dm_xml = zip_ref.read(data_model_target)
                                    dm_root = ET.fromstring(dm_xml)
                                    for pt in dm_root.findall('.//*[@val]'):
                                        val = pt.get('val')
                                        if val and not any(val in t for t in all_text):
                                            all_text.append(val)
                                except:
                                    pass
                
            return " ".join(all_text).strip()
        except Exception as e:
            print(f"Error extracting SmartArt text: {e}")
            return ""

def extract_text(pptx_file):  
    """Extract text from PowerPoint presentation with ultimate text extraction"""
    # Standard extraction first
    text_dict = {}
    slide_metadata = []
    extraction_stats = {
        "standard_elements": 0,
        "deep_elements": 0,
        "special_elements": 0
    }
    
    # Standard python-pptx extraction
    prs = Presentation(pptx_file)
    for index, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": index + 1,
            "title": "",
            "content": []
        }
        
        # Extract slide notes if any
        try:
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                for note_shape in slide.notes_slide.shapes:
                    if hasattr(note_shape, "text_frame") and note_shape.text_frame and note_shape.text.strip():
                        note_id = f"slide_{index+1}_notes"
                        note_text = note_shape.text.strip()
                        text_dict[note_id] = note_text
                        slide_info["content"].append(f"[Note: {note_text}]")
                        extraction_stats["standard_elements"] += 1
        except:
            pass
        
        # Process all shapes on the slide
        for shape_id, shape in enumerate(slide.shapes):
            # Basic text extraction
            if hasattr(shape, "text") and shape.text.strip():
                object_id = f"slide_{index+1}_shape_{shape_id}"
                text_dict[object_id] = shape.text.strip()
                slide_info["content"].append(shape.text.strip())
                extraction_stats["standard_elements"] += 1
                
                # Handle titles
                if hasattr(shape, "is_title") and shape.is_title:
                    slide_info["title"] = shape.text.strip()
                elif hasattr(shape, "is_placeholder") and shape.is_placeholder and hasattr(shape, "placeholder_format"):
                    try:
                        if shape.placeholder_format.type == 1:  # Title placeholder
                            slide_info["title"] = shape.text.strip()
                    except:
                        pass
            
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            cell_id = f"slide_{index+1}_table_{shape_id}_r{row_idx}_c{col_idx}"
                            text_dict[cell_id] = cell.text.strip()
                            slide_info["content"].append(cell.text.strip())
                            extraction_stats["standard_elements"] += 1
            
            # Try to get text from charts
            try:
                if hasattr(shape, "chart") and shape.chart:
                    # Get chart title
                    if hasattr(shape.chart, "chart_title") and shape.chart.chart_title and shape.chart.chart_title.text_frame:
                        chart_title_id = f"slide_{index+1}_chart_{shape_id}_title"
                        chart_title = shape.chart.chart_title.text_frame.text
                        text_dict[chart_title_id] = chart_title
                        slide_info["content"].append(chart_title)
                        extraction_stats["standard_elements"] += 1
            except:
                pass
        
        slide_metadata.append(slide_info)
    
    # Deep XML extraction for elements that python-pptx might miss
    with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
        # Process each slide
        slide_xmls = [f for f in zip_ref.namelist() if re.match(r'ppt/slides/slide[0-9]+\.xml', f)]
        
        for slide_xml in slide_xmls:
            # Extract slide number from filename
            slide_num = int(re.search(r'slide([0-9]+)\.xml', slide_xml).group(1))
            
            # Get slide relationships file
            slide_rels = slide_xml.replace('.xml', '.xml.rels')
            if slide_rels not in zip_ref.namelist():
                slide_rels = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
            
            # Parse slide XML
            try:
                slide_content = zip_ref.read(slide_xml)
                slide_root = ET.fromstring(slide_content)
            except:
                continue
            
            # Find all graphicData elements
            for graphic in slide_root.findall('.//a:graphicData', namespaces):
                uri = graphic.get('uri', '')
                
                # Process SmartArt
                if 'smartArt' in uri:
                    # Find the SmartArt relationship
                    for dgm in graphic.findall('.//a:dgm', namespaces):
                        if '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id' in dgm.attrib:
                            rel_id = dgm.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
                            smartart_text = extract_from_smartart(pptx_file, slide_rels, rel_id)
                            if smartart_text:
                                object_id = f"slide_{slide_num}_smartart_{len(text_dict)}"
                                text_dict[object_id] = smartart_text
                                extraction_stats["deep_elements"] += 1
                                
                                # Add to slide metadata
                                for slide_data in slide_metadata:
                                    if slide_data["slide_number"] == slide_num:
                                        slide_data["content"].append(smartart_text)
                                        break
                
                # Process WordArt and other graphics
                elif 'wordArt' in uri or 'diagram' in uri:
                    text = extract_text_from_element(graphic)
                    if text:
                        object_id = f"slide_{slide_num}_wordart_{len(text_dict)}"
                        text_dict[object_id] = text
                        extraction_stats["deep_elements"] += 1
                        
                        # Add to slide metadata
                        for slide_data in slide_metadata:
                            if slide_data["slide_number"] == slide_num:
                                slide_data["content"].append(text)
                                break
            
            # Find all text in the slide, including those that might be missed by python-pptx
            for para in slide_root.findall('.//a:p', namespaces):
                text = extract_text_from_element(para)
                # Only add text not already captured (avoid duplication)
                if text and not any(text in v for v in text_dict.values()):
                    object_id = f"slide_{slide_num}_text_{len(text_dict)}"
                    text_dict[object_id] = text
                    extraction_stats["deep_elements"] += 1
                    
            # Look for text in table cells (often missed)
            for tc in slide_root.findall('.//a:tc', namespaces):
                text = extract_text_from_element(tc)
                if text and not any(text in v for v in text_dict.values()):
                    object_id = f"slide_{slide_num}_table_cell_{len(text_dict)}"
                    text_dict[object_id] = text
                    extraction_stats["deep_elements"] += 1
                    
                    # Add to slide metadata
                    for slide_data in slide_metadata:
                        if slide_data["slide_number"] == slide_num:
                            slide_data["content"].append(text)
                            break
            
            # Look for text in comments
            for comment in slide_root.findall('.//p:cm', namespaces):
                text = extract_text_from_element(comment)
                if text and not any(text in v for v in text_dict.values()):
                    object_id = f"slide_{slide_num}_comment_{len(text_dict)}"
                    text_dict[object_id] = text
                    extraction_stats["deep_elements"] += 1
                    
                    # Add to slide metadata
                    for slide_data in slide_metadata:
                        if slide_data["slide_number"] == slide_num:
                            slide_data["content"].append(text)
                            break
                    
            # Comprehensive scanning for ALL slides - ensures we don't miss any text
            # Look for any text content in any element
            for elem in slide_root.findall('.//*'):
                if elem.text and elem.text.strip() and not any(elem.text.strip() in v for v in text_dict.values()):
                    special_id = f"slide_{slide_num}_special_{len(text_dict)}"
                    text_dict[special_id] = elem.text.strip()
                    extraction_stats["special_elements"] += 1
                    
                    # Add to slide metadata
                    for slide_data in slide_metadata:
                        if slide_data["slide_number"] == slide_num:
                            slide_data["content"].append(elem.text.strip())
                            break
                            
            # Handle the special case of alt text on images and shapes
            for elem in slide_root.findall('.//p:nvPr', namespaces):
                alt_text_elem = elem.find('.//a:altTxt', namespaces) or elem.find('.//p:extLst//p:ext//a14:alt', namespaces)
                if alt_text_elem is not None and alt_text_elem.text and alt_text_elem.text.strip():
                    if not any(alt_text_elem.text.strip() in v for v in text_dict.values()):
                        alt_id = f"slide_{slide_num}_alt_{len(text_dict)}"
                        text_dict[alt_id] = alt_text_elem.text.strip()
                        extraction_stats["special_elements"] += 1
                        
                        # Add to slide metadata
                        for slide_data in slide_metadata:
                            if slide_data["slide_number"] == slide_num:
                                slide_data["content"].append(alt_text_elem.text.strip())
                                break
    
    print(f"Ultimate extraction found {len(text_dict)} text elements")
    print(f"  - Standard extraction: {extraction_stats['standard_elements']} elements")
    print(f"  - Deep XML extraction: {extraction_stats['deep_elements']} elements")
    print(f"  - Special elements extraction: {extraction_stats['special_elements']} elements")
    
    return text_dict, slide_metadata

def update_slides(pptx_file, output_file, translated_texts):
    """Update PowerPoint presentation with translated text while preserving formatting"""
    prs = Presentation(pptx_file)
    updated_count = 0
    total_slides = len(prs.slides)
    
    # Process a text frame while preserving all formatting
    def update_text_frame(text_frame, new_text):
        nonlocal updated_count
        if not text_frame or not new_text:
            return False
        
        # Save formatting details from each paragraph
        formatting = []
        for para in text_frame.paragraphs:
            para_format = {
                "alignment": para.alignment if hasattr(para, "alignment") else None,
                "level": para.level if hasattr(para, "level") else None,
                "runs": []
            }
            
            # Save formatting for each run (text chunk with same formatting)
            for run in para.runs:
                run_format = {}
                if hasattr(run, "font") and run.font:
                    font = run.font
                    run_format["size"] = font.size if hasattr(font, "size") else None
                    run_format["bold"] = font.bold if hasattr(font, "bold") else None
                    run_format["italic"] = font.italic if hasattr(font, "italic") else None
                    run_format["underline"] = font.underline if hasattr(font, "underline") else None
                    
                    # Store font name and color as well
                    if hasattr(font, "name"):
                        run_format["name"] = font.name
                    if hasattr(font, "color") and hasattr(font.color, "rgb"):
                        run_format["color"] = font.color.rgb
                
                # Save hyperlink if present
                if hasattr(run, "hyperlink") and run.hyperlink:
                    run_format["hyperlink"] = run.hyperlink.address
                    
                para_format["runs"].append(run_format)
            
            formatting.append(para_format)
        
        # Clear and set the new text
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = new_text
        
        # Try to restore paragraph-level formatting for the first paragraph
        if formatting and formatting[0]:
            if formatting[0]["alignment"] is not None:
                p.alignment = formatting[0]["alignment"]
            if formatting[0]["level"] is not None:
                p.level = formatting[0]["level"]
        
        # Attempt to restore font formatting for the first run
        if formatting and formatting[0] and formatting[0]["runs"] and p.runs:
            run = p.runs[0]
            original_format = formatting[0]["runs"][0]
            
            if hasattr(run, "font"):
                if original_format.get("size") is not None:
                    run.font.size = original_format["size"]
                if original_format.get("bold") is not None:
                    run.font.bold = original_format["bold"]
                if original_format.get("italic") is not None:
                    run.font.italic = original_format["italic"]
                if original_format.get("underline") is not None:
                    run.font.underline = original_format["underline"]
                if original_format.get("name") is not None:
                    run.font.name = original_format["name"]
                if original_format.get("color") is not None:
                    run.font.color.rgb = original_format["color"]
                
                # Restore hyperlink if it was present
                if original_format.get("hyperlink") is not None and hasattr(run, "hyperlink"):
                    run.hyperlink.address = original_format["hyperlink"]
        
        updated_count += 1
        return True
    
    # Process a single shape and its text
    def process_shape(shape, shape_id, slide_idx):
        base_id = f"slide_{slide_idx+1}_shape_{shape_id}"
        updated = False
        
        # Update main text frame if present
        if base_id in translated_texts and hasattr(shape, "text_frame"):
            update_text_frame(shape.text_frame, translated_texts[base_id])
            updated = True
        
        # Handle charts
        try:
            if hasattr(shape, "chart") and shape.chart:
                # Update chart title
                chart_title_id = f"slide_{slide_idx+1}_chart_{shape_id}_title"
                if chart_title_id in translated_texts and hasattr(shape.chart, "chart_title") and hasattr(shape.chart.chart_title, "text_frame"):
                    update_text_frame(shape.chart.chart_title.text_frame, translated_texts[chart_title_id])
                    updated = True
        except:
            pass
        
        # Handle tables
        try:
            if hasattr(shape, "has_table") and shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_id = f"slide_{slide_idx+1}_table_{shape_id}_r{row_idx}_c{col_idx}"
                        if cell_id in translated_texts and hasattr(cell, "text_frame"):
                            update_text_frame(cell.text_frame, translated_texts[cell_id])
                            updated = True
        except:
            pass
                
        return updated
    
    # Update each slide with progress bar
    from tqdm import tqdm
    with tqdm(total=total_slides, desc="Updating slides", unit="slide") as pbar:
        for slide_idx, slide in enumerate(prs.slides):
            # Update slide notes if any
            try:
                if hasattr(slide, "notes_slide") and slide.notes_slide:
                    for note_shape in slide.notes_slide.shapes:
                        note_id = f"slide_{slide_idx+1}_notes"
                        if note_id in translated_texts and hasattr(note_shape, "text_frame"):
                            update_text_frame(note_shape.text_frame, translated_texts[note_id])
            except:
                pass
            
            # Update each shape on the slide
            for shape_idx, shape in enumerate(slide.shapes):
                process_shape(shape, shape_idx, slide_idx)
                
            # Update progress bar
            completion_percentage = int(100 * (slide_idx + 1) / total_slides)
            pbar.set_description(f"Updating slides: {completion_percentage}% complete")
            pbar.update(1)
            pbar.refresh()
    
    print(f"Updated {updated_count} text elements in the presentation")
    
    # Save the updated presentation to a new file
    prs.save(output_file)
    return output_file

# Include other functions from enhanced_pptx_translator.py (split_dict_into_smart_batches, repair_json, etc.)
def split_dict_into_smart_batches(input_dict, max_input_tokens=150000, prompt_tokens=2000):
    """
    Split a dictionary into batches based on estimated token count to optimize API usage.
    """
    # Function to estimate tokens in a string with better Unicode/multibyte handling
    def estimate_tokens(text):
        if text is None:
            return 0
            
        # Convert to string if not already
        text_str = str(text)
        
        # For Asian languages (CJK), use 1.5 characters per token as a conservative estimate
        cjk_chars = sum(1 for c in text_str if ord(c) > 0x4E00 and ord(c) < 0x9FFF)  # Basic CJK Unified Ideographs
        ascii_chars = len(text_str) - cjk_chars
        
        # 4 ASCII chars per token, ~1.5 CJK chars per token (rough estimate)
        token_estimate = (ascii_chars // 4) + (cjk_chars // 1.5) + 1  # Add 1 to round up
        
        return int(token_estimate)
    
    items = list(input_dict.items())
    batches = []
    current_batch = {}
    current_token_count = prompt_tokens
    
    # Sort items by estimated token length (optional)
    items.sort(key=lambda x: estimate_tokens(x[1]), reverse=True)
    
    for key, value in items:
        item_tokens = estimate_tokens(key) + estimate_tokens(value) + 10  # +10 for JSON formatting
        
        if current_token_count + item_tokens > max_input_tokens and current_batch:
            batches.append(current_batch)
            current_batch = {}
            current_token_count = prompt_tokens
        
        current_batch[key] = value
        current_token_count += item_tokens
    
    if current_batch:
        batches.append(current_batch)
    
    total_items = len(input_dict)
    batch_sizes = [len(batch) for batch in batches]
    avg_batch_size = sum(batch_sizes) / len(batches) if batches else 0
    
    print(f"Created {len(batches)} batches from {total_items} items")
    print(f"Batch sizes: min={min(batch_sizes) if batches else 0}, max={max(batch_sizes) if batches else 0}, avg={avg_batch_size:.1f}")
    print(f"Estimated token usage efficiency: {(sum(batch_sizes)/total_items)*100:.1f}%")
    
    return batches

def repair_json(json_content):
    """
    More robust JSON repair function that can handle various common issues including Unicode.
    """
    original_content = json_content
    try:
        return json.loads(json_content)
    except json.JSONDecodeError as e:
        print(f"Initial JSON parsing error: {e}")
        
        # Try to handle Unicode/multibyte character issues
        try:
            # Ensure the content is properly encoded
            if isinstance(json_content, bytes):
                json_content = json_content.decode('utf-8', errors='replace')
            
            # Remove or replace invisible/control characters that might cause issues
            json_content = re.sub(r'[\x00-\x1F\x7F]', '', json_content)
        except Exception as enc_err:
            print(f"Error handling encoding: {enc_err}")
        
        if "Unterminated string" in str(e):
            error_info = str(e)
            line_match = re.search(r'line (\d+)', error_info)
            col_match = re.search(r'column (\d+)', error_info)
            
            if line_match and col_match:
                line_num = int(line_match.group(1))
                col_num = int(col_match.group(1))
                lines = json_content.split('\n')
                if 0 <= line_num-1 < len(lines):
                    line = lines[line_num-1]
                    if col_num-1 < len(line) and line[col_num-1] == '"':
                        line = line[:col_num-1] + '\\"' + line[col_num:]
                    else:
                        line = line + '"'
                    lines[line_num-1] = line
                    json_content = '\n'.join(lines)
        
        # Handle improperly escaped characters and unicode escapes
        json_content = re.sub(r'(?<!\\)\\(?!["\\/bfnrt]|u[0-9a-fA-F]{4})', r'\\\\', json_content)
        
        # Fix unbalanced braces
        brace_count = json_content.count('{') - json_content.count('}')
        if brace_count > 0:
            json_content = json_content + ('}' * brace_count)
        elif brace_count < 0:
            for _ in range(-brace_count):
                json_content = json_content.rstrip().rstrip('}').rstrip()
        
        # Fix common JSON syntax issues
        json_content = re.sub(r',\s*}', '}', json_content)
        json_content = re.sub(r',\s*]', ']', json_content)
        
        # Ensure property names are properly quoted
        def fix_property_names(match):
            prop = match.group(1)
            if not (prop.startswith('"') and prop.endswith('"')):
                return f'"{prop}":'
            return match.group(0)
        
        json_content = re.sub(r'([a-zA-Z0-9_]+):', fix_property_names, json_content)
        
        # Try to parse the repaired JSON
        try:
            return json.loads(json_content)
        except json.JSONDecodeError as e2:
            print(f"JSON repair attempt failed: {e2}")
            
            # Fallback: extract key-value pairs using regex
            result = {}
            # Modified pattern to handle Unicode characters better
            pattern = r'"((?:[^"\\]|\\.)*)"\s*:\s*"((?:[^"\\]|\\.)*)"'
            for match in re.finditer(pattern, original_content):
                try:
                    key, value = match.groups()
                    # Unescape escaped quotes in the extracted strings
                    key = key.replace('\\"', '"')
                    value = value.replace('\\"', '"')
                    result[key] = value
                except Exception as extract_err:
                    print(f"Error extracting key-value pair: {extract_err}")
            
            # Also try to capture numeric values
            num_pattern = r'"((?:[^"\\]|\\.)*)"\s*:\s*([0-9]+(?:\.[0-9]+)?)'
            for match in re.finditer(num_pattern, original_content):
                try:
                    key, value = match.groups()
                    key = key.replace('\\"', '"')
                    # Convert to int or float as appropriate
                    if '.' in value:
                        result[key] = float(value)
                    else:
                        result[key] = int(value)
                except Exception as extract_err:
                    print(f"Error extracting numeric value: {extract_err}")
            
            if result:
                print(f"Managed to extract {len(result)} key-value pairs through regex")
                return result
            
            raise e

def extract_json_blocks(text):
    """
    Extract valid JSON blocks from text that might contain multiple partial JSON objects.
    Enhanced to better handle Unicode and Japanese characters.
    """
    # First try to find complete JSON objects
    try:
        # Clean the text: remove any leading/trailing non-JSON content
        text = re.sub(r'^[^{]*', '', text)  # Remove anything before the first {
        text = re.sub(r'[^}]*$', '', text)  # Remove anything after the last }
        
        # Try extracting json blocks with a more robust pattern
        # This pattern tries to match balanced { } pairs
        depth = 0
        start = -1
        blocks = []
        
        for i, char in enumerate(text):
            if char == '{':
                if depth == 0:
                    start = i
                depth += 1
            elif char == '}':
                depth -= 1
                if depth == 0 and start != -1:
                    blocks.append(text[start:i+1])
                    start = -1
        
        if not blocks:
            # Fallback to simpler regex if the balanced matching didn't work
            potential_blocks = re.findall(r'({[^{]*?})', text)
            blocks = potential_blocks
        
        valid_blocks = []
        for block in blocks:
            try:
                # Fix common issues that might occur in the JSON block
                block = re.sub(r',\s*}', '}', block)  # Remove trailing commas
                block = re.sub(r'([a-zA-Z0-9_]+):', r'"\1":', block)  # Quote unquoted keys
                
                parsed = json.loads(block)
                valid_blocks.append(parsed)
            except json.JSONDecodeError:
                try:
                    # Try to repair the JSON before giving up
                    repaired = repair_json(block)
                    if repaired:
                        valid_blocks.append(repaired)
                except:
                    pass
        
        if valid_blocks:
            combined = {}
            for block in valid_blocks:
                combined.update(block)
            return combined
            
    except Exception as e:
        print(f"Error in JSON block extraction: {e}")
        
    # Fallback: just try a direct key-value extraction
    try:
        result = {}
        # Look for key-value pairs directly, handling Unicode properly
        pattern = r'"((?:[^"\\]|\\.)*)"\s*:\s*"((?:[^"\\]|\\.)*)"'
        for match in re.finditer(pattern, text):
            key, value = match.groups()
            # Unescape escaped quotes
            key = key.replace('\\"', '"')
            value = value.replace('\\"', '"')
            result[key] = value
        
        if result:
            print(f"Direct key-value extraction found {len(result)} pairs")
            return result
    except Exception as e:
        print(f"Fallback extraction error: {e}")
        
    return None

def setup_recovery_system(file_id, text_dict, slide_metadata, source_language, target_language, resume_file=None):
    """
    Set up a recovery system for batch processing.
    """
    recovery_dir = "translation_recovery"
    os.makedirs(recovery_dir, exist_ok=True)
    
    if resume_file and os.path.exists(resume_file):
        with open(resume_file, 'r', encoding='utf-8') as f:
            recovery_state = json.load(f)
        print(f"Resuming translation from recovery file: {resume_file}")
        recovery_file = resume_file
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        recovery_file = os.path.join(recovery_dir, f"recovery_ultimate_{file_id}_{timestamp}.json")
        recovery_state = {
            "file_id": file_id,
            "completed_batches": [],
            "failed_batches": [],
            "translated_items": {},
            "source_language": source_language,
            "target_language": target_language,
            "total_items": len(text_dict),
            "start_time": timestamp,
            "last_updated": timestamp
        }
        with open(recovery_file, 'w', encoding='utf-8') as f:
            json.dump(recovery_state, f, ensure_ascii=False, indent=2)
        print(f"Created new recovery file: {recovery_file}")
    
    def save_recovery_state():
        recovery_state["last_updated"] = datetime.now().strftime("%Y%m%d_%H%M%S")
        with open(recovery_file, 'w', encoding='utf-8') as f:
            json.dump(recovery_state, f, ensure_ascii=False, indent=2)
    
    return recovery_state, recovery_file, save_recovery_state

def translate_batch(batch, batch_index, slide_metadata, source_language, target_language, api_key=None, max_retries=3, cost_tracker=None):
    """
    Translate a single batch with retry logic.
    """
    batch_copy = batch.copy()
    
    def clean_text(text):
        return text.replace('\\n', '\n').replace('\\u000b', '\v').replace('\\t', '\t')
        
    # Calculate estimated cost based on prompt and completion tokens
    def estimate_cost(prompt_tokens, completion_tokens, model="claude-3-7-sonnet"):
        # Claude 3.5 Sonnet pricing: $3 per 1M input tokens, $15 per 1M output tokens
        # Claude 3 Sonnet pricing: $3 per 1M input tokens, $15 per 1M output tokens
        rates = {
            "claude-3-7-sonnet": {"input": 3.0 / 1000000, "output": 15.0 / 1000000},
            "claude-3-sonnet": {"input": 3.0 / 1000000, "output": 15.0 / 1000000}, 
            "claude-3-opus": {"input": 15.0 / 1000000, "output": 75.0 / 1000000},
            "claude-3-haiku": {"input": 0.25 / 1000000, "output": 1.25 / 1000000}
        }
        
        model_rates = rates.get(model, rates["claude-3-7-sonnet"])  # Default to sonnet if not found
        input_cost = prompt_tokens * model_rates["input"]
        output_cost = completion_tokens * model_rates["output"]
        total_cost = input_cost + output_cost
        
        return {
            "input_tokens": prompt_tokens,
            "output_tokens": completion_tokens,
            "input_cost": input_cost,
            "output_cost": output_cost,
            "total_cost": total_cost
        }
    
    # Use provided API key or fall back to environment variable
    api_key = api_key or os.getenv("CLAUDE_API_KEY")
    if not api_key:
        raise ValueError("Claude API key must be provided either as an argument or via CLAUDE_API_KEY environment variable")
        
    client = anthropic.Anthropic(
        api_key=api_key,
        default_headers={
            "anthropic-beta": "output-128k-2025-02-19"
        }
    )
    
    structured_context = json.dumps(slide_metadata, ensure_ascii=False, indent=2)
    
    system_prompt = f"""You are a professional translator. Translate from {source_language} to {target_language}.
Ensure consistency in terminology and contextual meaning.

IMPORTANT: If you encounter text that appears to already be in {target_language}, preserve it exactly as is without any changes.
Do not translate text that is already in {target_language}.

PRIVACY NOTICE: Do not store, learn from, or retain any of the content provided for translation.
This is confidential material that should only be processed for immediate translation purposes."""
    
    user_message = f"""
Translate the following JSON object from {source_language} to {target_language}.
Consider the structured slide context provided below for context and consistency.

IMPORTANT INSTRUCTIONS:
- If any text appears to already be in {target_language}, keep it exactly as is.
- Only translate text that is in {source_language}.
- Do NOT include escape sequences for newlines (\\n) or other characters - use the actual characters.
- Return VALID JSON format with all keys and values properly enclosed in double quotes.
- Ensure all property names and string values are properly quoted with double quotes.
- Do not include any trailing commas.

This is batch {batch_index} with {len(batch_copy)} items.

Slide Context:
{structured_context}

Now translate the following structured JSON object while preserving its format:
{json.dumps(batch_copy, ensure_ascii=False, indent=2)}

Reply ONLY with the translated JSON. The JSON MUST be valid and parseable.
"""

    for retry in range(max_retries + 1):
        try:
            response = client.messages.create(
                model="claude-3-7-sonnet-20250219",
                system=system_prompt,
                max_tokens=4000,
                messages=[
                    {"role": "user", "content": user_message}
                ],
                metadata={
                    "user_id": "anonymous_user"
                }
            )

            # Track token usage and cost
            prompt_tokens = response.usage.input_tokens
            completion_tokens = response.usage.output_tokens
            
            batch_cost = estimate_cost(prompt_tokens, completion_tokens, "claude-3-7-sonnet")
            
            # Update the cost tracker if provided
            if cost_tracker is not None:
                cost_tracker["total_input_tokens"] += prompt_tokens
                cost_tracker["total_output_tokens"] += completion_tokens
                cost_tracker["total_input_cost"] += batch_cost["input_cost"]
                cost_tracker["total_output_cost"] += batch_cost["output_cost"]
                cost_tracker["total_cost"] += batch_cost["total_cost"]
                cost_tracker["api_calls"] += 1
            
            print(f"Batch {batch_index} token usage: {prompt_tokens} input + {completion_tokens} output tokens")
            print(f"Batch {batch_index} cost: ${batch_cost['total_cost']:.4f} (${batch_cost['input_cost']:.4f} input + ${batch_cost['output_cost']:.4f} output)")
            
            # Show running total cost
            if cost_tracker is not None:
                print(f"Running total: ${cost_tracker['total_cost']:.4f} for {cost_tracker['api_calls']} API calls")
            
            translated_text = response.content[0].text
            
            if "```json" in translated_text:
                json_content = translated_text.split("```json")[1].split("```")[0].strip()
            elif "```" in translated_text:
                json_content = translated_text.split("```")[1].strip()
            else:
                json_content = translated_text.strip()
            
            try:
                batch_result = json.loads(json_content)
            except json.JSONDecodeError as e:
                try:
                    batch_result = repair_json(json_content)
                except Exception as e2:
                    extracted_result = extract_json_blocks(json_content)
                    if extracted_result:
                        batch_result = extracted_result
                        print(f"Extracted {len(batch_result)} items through JSON block extraction")
                    else:
                        if retry < max_retries:
                            print(f"JSON parsing failed on attempt {retry+1}, retrying...")
                            time.sleep(3)
                            continue
                        else:
                            raise e
            
            for key, value in batch_result.items():
                if isinstance(value, str):
                    batch_result[key] = clean_text(value)
            
            print(f"Successfully processed batch {batch_index}")
            return batch_result
                
        except Exception as e:
            if retry < max_retries:
                print(f"Error in batch {batch_index} (attempt {retry+1}): {e}")
                print(f"Retrying in 5 seconds...")
                time.sleep(5)
            else:
                print(f"All {max_retries + 1} attempts failed for batch {batch_index}: {e}")
                raise e

def translate_text(text_dict, slide_metadata, source_language, target_language, resume_file=None, api_key=None):
    # Use provided API key or fall back to environment variable
    api_key = api_key or os.getenv("CLAUDE_API_KEY")
    if not api_key:
        raise ValueError("Claude API key must be provided either as an argument or via CLAUDE_API_KEY environment variable")
        
    client = anthropic.Anthropic(
        api_key=api_key,
        default_headers={
            "anthropic-beta": "output-128k-2025-02-19"
        }
    )
    
    # Initialize cost tracking
    cost_tracker = {
        "total_input_tokens": 0,
        "total_output_tokens": 0,
        "total_input_cost": 0.0,
        "total_output_cost": 0.0,
        "total_cost": 0.0,
        "api_calls": 0
    }
    
    def deduplicate_content(input_dict):
        content_to_keys = {}
        for key, value in input_dict.items():
            if value in content_to_keys:
                content_to_keys[value].append(key)
            else:
                content_to_keys[value] = [key]
        
        unique_content = {}
        duplicates_map = {}
        
        for content, keys in content_to_keys.items():
            representative_key = keys[0]
            unique_content[representative_key] = content
            for key in keys:
                duplicates_map[key] = representative_key
        
        print(f"Found {len(input_dict) - len(unique_content)} duplicate content items")
        print(f"Reduced from {len(input_dict)} to {len(unique_content)} unique content items to translate")
        
        return unique_content, duplicates_map
    
    # Use a simple file ID based on timestamp if not provided
    file_id = os.path.basename(str(time.time()).replace('.', ''))
    
    recovery_state, recovery_file, save_recovery_state = setup_recovery_system(
        file_id, text_dict, slide_metadata, source_language, target_language, resume_file
    )
    
    if not recovery_state["translated_items"]:
        unique_text_dict, duplicates_map = deduplicate_content(text_dict)
        recovery_state["duplicates_map"] = duplicates_map
        save_recovery_state()
    else:
        unique_text_dict = {k: text_dict[k] for k in text_dict.keys() 
                          if k not in [rep_key for rep_key in recovery_state["duplicates_map"].values()]}
        duplicates_map = recovery_state["duplicates_map"]
        print(f"Resumed with {len(recovery_state['translated_items'])} already translated items")
    
    remaining_dict = {k: v for k, v in unique_text_dict.items() 
                     if k not in recovery_state["translated_items"]}
    
    if not remaining_dict:
        print("All items have already been translated. Nothing to do.")
        full_translated_dict = recovery_state["translated_items"].copy()
    else:
        # For Japanese or other multibyte languages, use smaller batches
        max_tokens = 50000 if target_language in ["ja", "zh", "ko"] else 100000
        prompt_tokens = 2000
        
        print(f"Using smaller batch size for {target_language} translation" if target_language in ["ja", "zh", "ko"] else "Using standard batch size")
        batches = split_dict_into_smart_batches(remaining_dict, max_input_tokens=max_tokens, prompt_tokens=prompt_tokens)
        print(f"Splitting translation into {len(batches)} batches")
        
        unique_translated_dict = recovery_state["translated_items"].copy()
        
        with tqdm(total=len(batches), desc="Translating", unit="batch") as pbar:
            for batch_index, batch in enumerate(batches):
                batch_id = f"batch_{batch_index+1}"
                if batch_id in recovery_state["completed_batches"]:
                    print(f"Skipping already completed batch {batch_id}")
                    pbar.update(1)
                    continue
                
                print(f"\nProcessing batch {batch_index+1} of {len(batches)} with {len(batch)} items...")
                
                try:
                    batch_result = translate_batch(
                        batch, batch_index+1, slide_metadata, 
                        source_language, target_language, api_key=api_key, 
                        cost_tracker=cost_tracker
                    )
                    
                    unique_translated_dict.update(batch_result)
                    recovery_state["translated_items"].update(batch_result)
                    recovery_state["completed_batches"].append(batch_id)
                    save_recovery_state()
                    
                except Exception as e:
                    print(f"Error in batch {batch_index+1}: {e}")
                    recovery_state["failed_batches"].append({
                        "batch_id": batch_id,
                        "keys": list(batch.keys()),
                        "error": str(e)
                    })
                    save_recovery_state()
                    print("Continuing with next batch...")
                
                pbar.update(1)
                completion_percentage = int(100 * (batch_index + 1) / len(batches))
                pbar.set_description(f"Translating: {completion_percentage}% complete")
                # Force refresh the progress bar display
                pbar.refresh()
        
        print(f"\nTranslation of unique content completed with {len(unique_translated_dict)} items out of {len(unique_text_dict)} unique items")
        
        if recovery_state["failed_batches"]:
            print(f"\nRetrying {len(recovery_state['failed_batches'])} failed batches with smaller chunks...")
            
            for failed_batch in list(recovery_state["failed_batches"]):
                batch_id = failed_batch["batch_id"]
                keys = failed_batch["keys"]
                
                retry_batch = {k: text_dict[k] for k in keys if k in text_dict}
                # Use even smaller chunks for Japanese
                chunk_size = max(3, len(retry_batch) // 20 if target_language in ["ja", "zh", "ko"] else len(retry_batch) // 10)
                retry_items = list(retry_batch.items())
                sub_batches = [dict(retry_items[i:i+chunk_size]) 
                               for i in range(0, len(retry_items), chunk_size)]
                
                print(f"Split failed batch {batch_id} into {len(sub_batches)} smaller chunks of size ~{chunk_size}")
                
                for i, sub_batch in enumerate(sub_batches):
                    sub_id = f"{batch_id}_sub_{i+1}"
                    
                    try:
                        print(f"Processing sub-batch {i+1}/{len(sub_batches)} for failed batch {batch_id}")
                        sub_result = translate_batch(
                            sub_batch, f"{batch_id}.{i+1}", slide_metadata, 
                            source_language, target_language, api_key=api_key, 
                            max_retries=3, cost_tracker=cost_tracker
                        )
                        unique_translated_dict.update(sub_result)
                        recovery_state["translated_items"].update(sub_result)
                        recovery_state["completed_batches"].append(sub_id)
                        save_recovery_state()
                        
                    except Exception as e:
                        print(f"Error in sub-batch {i+1} of failed batch {batch_id}: {e}")
                        continue
                
                recovery_state["failed_batches"].remove(failed_batch)
                save_recovery_state()
        
        full_translated_dict = {}
        for key, value in unique_translated_dict.items():
            full_translated_dict[key] = value
        
        for original_key, rep_key in duplicates_map.items():
            if rep_key in unique_translated_dict and original_key != rep_key:
                full_translated_dict[original_key] = unique_translated_dict[rep_key]
    
    print(f"Reconstructed full translation dictionary with {len(full_translated_dict)} items")
    
    missing_keys = set(text_dict.keys()) - set(full_translated_dict.keys())
    if missing_keys:
        print(f"Warning: {len(missing_keys)} keys were not translated: {list(missing_keys)[:5]}...")
        if len(missing_keys) > 0:
            print(f"Attempting to translate {len(missing_keys)} missing keys in a final batch...")
            missing_dict = {k: text_dict[k] for k in missing_keys if k in text_dict}
            
            # Process final batch in smaller chunks for CJK languages
            final_chunk_size = 50 if target_language in ["ja", "zh", "ko"] else 100
            missing_items = list(missing_dict.items())
            
            # Break final batch into smaller manageable chunks
            final_chunk_count = (len(missing_items) + final_chunk_size - 1) // final_chunk_size  # Ceiling division
            
            if final_chunk_count > 1:
                print(f"Processing final {len(missing_items)} items in {final_chunk_count} smaller chunks")
                
                for chunk_idx in range(final_chunk_count):
                    start_idx = chunk_idx * final_chunk_size
                    end_idx = min(start_idx + final_chunk_size, len(missing_items))
                    chunk_dict = dict(missing_items[start_idx:end_idx])
                    
                    try:
                        print(f"Processing final chunk {chunk_idx+1}/{final_chunk_count} with {len(chunk_dict)} items")
                        
                        # Use the translate_batch function for consistent handling
                        chunk_result = translate_batch(
                            chunk_dict, f"final_{chunk_idx+1}", slide_metadata,
                            source_language, target_language, api_key=api_key,
                            max_retries=3, cost_tracker=cost_tracker
                        )
                        
                        full_translated_dict.update(chunk_result)
                        recovery_state["translated_items"].update(chunk_result)
                        save_recovery_state()
                        
                    except Exception as e:
                        print(f"Error processing final chunk {chunk_idx+1}: {e}")
            else:
                # Single final batch
                try:
                    # Use the translate_batch function for consistency
                    final_batch = translate_batch(
                        missing_dict, "final", slide_metadata,
                        source_language, target_language, api_key=api_key,
                        max_retries=3, cost_tracker=cost_tracker
                    )
                    
                    full_translated_dict.update(final_batch)
                    recovery_state["translated_items"].update(final_batch)
                    save_recovery_state()
                    
                except Exception as e:
                    print(f"Failed to process final batch: {e}")
            
            missing_keys = set(text_dict.keys()) - set(full_translated_dict.keys())
            if missing_keys:
                print(f"Final warning: {len(missing_keys)} keys still not translated: {list(missing_keys)[:5]}...")
            else:
                print("All items successfully translated!")
    else:
        print("All items successfully translated!")
    
    # Print cost summary
    print("\n=== API Cost Summary ===")
    print(f"Total API calls: {cost_tracker['api_calls']}")
    print(f"Total input tokens: {cost_tracker['total_input_tokens']:,}")
    print(f"Total output tokens: {cost_tracker['total_output_tokens']:,}")
    print(f"Total tokens: {cost_tracker['total_input_tokens'] + cost_tracker['total_output_tokens']:,}")
    print(f"Input cost: ${cost_tracker['total_input_cost']:.4f}")
    print(f"Output cost: ${cost_tracker['total_output_cost']:.4f}")
    print(f"Total cost: ${cost_tracker['total_cost']:.4f}")
    
    return full_translated_dict
    
def list_recovery_files():
    """List all available recovery files and their status"""
    recovery_dir = "translation_recovery"
    if not os.path.exists(recovery_dir):
        print("No recovery directory found.")
        return
    
    recovery_files = [f for f in os.listdir(recovery_dir) if f.endswith(".json")]
    
    if not recovery_files:
        print("No recovery files found.")
        return
    
    print(f"Found {len(recovery_files)} recovery files:")
    for f in recovery_files:
        file_path = os.path.join(recovery_dir, f)
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                total = data.get("total_items", 0)
                translated = len(data.get("translated_items", {}))
                failed = len(data.get("failed_batches", []))
                progress = (translated / total * 100) if total > 0 else 0
                
                print(f"  {f}")
                print(f"    Progress: {progress:.1f}% ({translated}/{total} items)")
                print(f"    Failed batches: {failed}")
                print(f"    Start time: {data.get('start_time', 'unknown')}")
                print(f"    Last updated: {data.get('last_updated', 'unknown')}")
                print()
        except Exception as e:
            print(f"  {f} - Error reading file: {e}")

def translate_pptx(input_file, output_file, source_language="en", target_language="fr", resume_file=None, api_key=None):
    """Main function to translate PowerPoint files with ultimate text extraction"""
    print(f"Extracting text from {input_file} with ultimate extraction...")
    text_dict, slide_metadata = extract_text(input_file)
    print(f"Found {len(text_dict)} text elements across {len(slide_metadata)} slides")
    
    print(f"Translating from {source_language} to {target_language}...")
    translated_texts = translate_text(text_dict, slide_metadata, source_language, target_language, resume_file, api_key=api_key)
    
    print(f"Updating PowerPoint with translated text while preserving formatting...")
    update_slides(input_file, output_file, translated_texts)
    
    print(f"Translation completed!")
    print(f"Translated presentation saved as: {output_file}")
    
    return output_file

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Ultimate PowerPoint PPTX Translator")
    parser.add_argument("--resume", help="Resume translation from a recovery file")
    parser.add_argument("--list-recovery", action="store_true", help="List available recovery files")
    parser.add_argument("--input-file", help="Input PowerPoint file (.pptx)")
    parser.add_argument("--output-file", help="Output PowerPoint file (.pptx)")
    parser.add_argument("--source-language", help="Source language code (e.g., en)")
    parser.add_argument("--target-language", help="Target language code (e.g., ja)")
    parser.add_argument("--api-key", help="Claude API Key (can also be set as CLAUDE_API_KEY environment variable)")
    
    args = parser.parse_args()
    
    if args.list_recovery:
        list_recovery_files()
        sys.exit(0)
    
    if args.input_file:
        input_file = args.input_file
    else:
        input_file = input("Enter path to PowerPoint file (.pptx): ")
    
    if not os.path.exists(input_file):
        print(f"Error: Input file '{input_file}' does not exist")
        sys.exit(1)
    
    if args.output_file:
        output_file = args.output_file
    else:
        # Create a default output filename based on the input filename
        base_name = os.path.splitext(os.path.basename(input_file))[0]
        if args.target_language:
            output_file = f"{base_name}_{args.target_language}_ultimate.pptx"
        else:
            output_file = f"{base_name}_translated_ultimate.pptx"
    
    if args.source_language:
        source_language = args.source_language
    else:
        source_language = input("Enter source language (e.g., en for English): ")
    
    if args.target_language:
        target_language = args.target_language
    else:
        target_language = input("Enter target language (e.g., fr for French): ")
    
    # Run the translation with all the provided parameters
    translate_pptx(input_file, output_file, source_language, target_language, args.resume, api_key=args.api_key)