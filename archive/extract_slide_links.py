#!/usr/bin/env python3
import os
import sys
from pptx import Presentation
import xml.etree.ElementTree as ET
import zipfile

def extract_slide_links(pptx_file, slide_numbers=None):
    """Extract links from specific slides in a PowerPoint presentation"""
    
    if not os.path.exists(pptx_file):
        print(f"Error: File '{pptx_file}' not found.")
        return None
    
    prs = Presentation(pptx_file)
    
    # If no slide numbers provided, process all slides
    if slide_numbers is None:
        slide_numbers = list(range(1, len(prs.slides) + 1))
    
    results = {}
    
    # Extract links from python-pptx accessible elements
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        
        if slide_num not in slide_numbers:
            continue
            
        slide_links = []
        
        # Process shapes that might have hyperlinks
        for shape in slide.shapes:
            # Text hyperlinks
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run, "hyperlink") and run.hyperlink and run.hyperlink.address:
                            slide_links.append({
                                "text": run.text,
                                "url": run.hyperlink.address,
                                "type": "text_link"
                            })
            
            # Shape hyperlinks
            if hasattr(shape, "click_action") and shape.click_action:
                if hasattr(shape.click_action, "hyperlink") and shape.click_action.hyperlink.address:
                    link_text = shape.text if hasattr(shape, "text") else "Shape link"
                    slide_links.append({
                        "text": link_text,
                        "url": shape.click_action.hyperlink.address,
                        "type": "shape_link"
                    })
        
        # Now use direct XML access to find links that python-pptx might miss
        try:
            with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
                # Get slide's XML content
                slide_xml_path = f"ppt/slides/slide{slide_num}.xml"
                if slide_xml_path in zip_ref.namelist():
                    slide_xml = zip_ref.read(slide_xml_path)
                    root = ET.fromstring(slide_xml)
                    
                    # XML namespaces in PPTX files
                    namespaces = {
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                    }
                    
                    # Get slide relationship XML (contains actual URLs)
                    rels_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
                    if rels_path in zip_ref.namelist():
                        rels_xml = zip_ref.read(rels_path)
                        rels_root = ET.fromstring(rels_xml)
                        
                        # Create a mapping of relationship IDs to targets (URLs)
                        rel_targets = {}
                        for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                            rel_id = rel.get('Id')
                            rel_target = rel.get('Target')
                            rel_type = rel.get('Type')
                            
                            # Only interested in hyperlinks
                            if 'hyperlink' in rel_type:
                                rel_targets[rel_id] = rel_target
                        
                        # Find all hyperlink references in the slide XML
                        for link_elem in root.findall('.//a:hlinkClick', namespaces):
                            rel_id = link_elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                            if rel_id in rel_targets:
                                # Try to find the text associated with this link
                                parent = link_elem.getparent().getparent()
                                link_text = "Unknown text"
                                
                                # Look for text in this element or its children
                                text_elems = parent.findall('.//a:t', namespaces)
                                if text_elems:
                                    link_text = ' '.join(t.text for t in text_elems if t.text)
                                
                                slide_links.append({
                                    "text": link_text,
                                    "url": rel_targets[rel_id],
                                    "type": "xml_link"
                                })
        except Exception as e:
            print(f"Error processing slide {slide_num} XML: {e}")
            
        # Store the results for this slide
        if slide_links:
            results[slide_num] = slide_links
    
    return results

def display_links(links_dict):
    """Display the extracted links in a readable format"""
    if not links_dict:
        print("No links found in the specified slides.")
        return
    
    total_links = sum(len(links) for links in links_dict.values())
    print(f"Found {total_links} links across {len(links_dict)} slides:")
    print()
    
    for slide_num, links in sorted(links_dict.items()):
        print(f"Slide {slide_num}: {len(links)} links")
        for i, link in enumerate(links, 1):
            print(f"  {i}. \"{link['text']}\" -> {link['url']} ({link['type']})")
        print()

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_slide_links.py <pptx_file> [slide_numbers]")
        print("Example: python extract_slide_links.py presentation.pptx 1 5 20")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    
    # If slide numbers are provided, convert them to integers
    slide_numbers = None
    if len(sys.argv) > 2:
        try:
            slide_numbers = [int(x) for x in sys.argv[2:]]
        except ValueError:
            print("Error: Slide numbers must be integers")
            sys.exit(1)
    
    links = extract_slide_links(pptx_file, slide_numbers)
    display_links(links)