#!/usr/bin/env python3
import zipfile
import xml.etree.ElementTree as ET
import os
import re
import json
from pptx import Presentation
import argparse

# XML namespaces used in PPTX files
namespaces = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006'
}

def extract_text_from_element(element):
    """Extract text from an XML element and its children."""
    text = ""
    
    # Extract text from a:t elements (text runs)
    for t in element.findall('.//a:t', namespaces):
        if t.text:
            text += t.text + " "
    
    return text.strip()

def extract_from_smartart(pptx_file, rels_path, rel_id):
    """Extract text from SmartArt diagrams using direct XML processing."""
    with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
        # Find the relationship target for the SmartArt
        rels_xml = zip_ref.read(rels_path)
        rels_root = ET.fromstring(rels_xml)
        
        # Find the target for this relationship ID
        target = None
        for rel in rels_root.findall('.//Relationship', namespaces):
            if rel.get('Id') == rel_id and rel.get('Type').endswith('diagramData'):
                target = rel.get('Target')
                break
        
        if not target:
            return ""
        
        # Convert the target path to the correct format
        if target.startswith('../'):
            target = target.replace('../', '')
        else:
            slide_dir = os.path.dirname(rels_path)
            target = os.path.join(os.path.dirname(slide_dir), target)
        
        # Read the diagram data
        try:
            diagram_xml = zip_ref.read(target)
            diagram_root = ET.fromstring(diagram_xml)
            
            # Extract text from text elements in the diagram
            text = ""
            for t_element in diagram_root.findall('.//a:t', namespaces):
                if t_element.text:
                    text += t_element.text + " "
            
            return text.strip()
        except Exception as e:
            print(f"Error extracting SmartArt text: {e}")
            return ""

def deep_extract_text(pptx_file, output_json=None):
    """Extract text from PowerPoint file using both python-pptx and direct XML processing."""
    # Standard extraction first
    text_dict = {}
    slide_metadata = []
    
    # Use python-pptx to extract what it can
    prs = Presentation(pptx_file)
    for index, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": index + 1,
            "title": "",
            "content": []
        }
        
        # Extract text from each shape using python-pptx
        for shape_id, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                object_id = f"slide_{index+1}_shape_{shape_id}"
                text_dict[object_id] = shape.text.strip()
                slide_info["content"].append(shape.text.strip())
                
                # Check if this is a title shape
                if hasattr(shape, "is_title") and shape.is_title:
                    slide_info["title"] = shape.text.strip()
            
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            cell_id = f"slide_{index+1}_table_{shape_id}_r{row_idx}_c{col_idx}"
                            text_dict[cell_id] = cell.text.strip()
                            slide_info["content"].append(cell.text.strip())
        
        slide_metadata.append(slide_info)
    
    # Now let's do deep XML processing for elements that python-pptx might miss
    with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
        # Process each slide
        slide_xmls = [f for f in zip_ref.namelist() if re.match(r'ppt/slides/slide[0-9]+\.xml', f)]
        
        for slide_xml in slide_xmls:
            # Extract slide number from filename
            slide_num = int(re.search(r'slide([0-9]+)\.xml', slide_xml).group(1))
            
            # Get slide relationships file
            slide_rels = slide_xml.replace('.xml', '.xml.rels')
            if slide_rels not in zip_ref.namelist():
                slide_rels = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
            
            # Parse slide XML
            slide_content = zip_ref.read(slide_xml)
            slide_root = ET.fromstring(slide_content)
            
            # Find all graphicData elements
            for graphic in slide_root.findall('.//a:graphicData', namespaces):
                uri = graphic.get('uri', '')
                
                # Process SmartArt
                if 'smartArt' in uri:
                    # Find the SmartArt relationship
                    for dgm in graphic.findall('.//a:dgm', namespaces):
                        if '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id' in dgm.attrib:
                            rel_id = dgm.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
                            smartart_text = extract_from_smartart(pptx_file, slide_rels, rel_id)
                            if smartart_text:
                                object_id = f"slide_{slide_num}_smartart_{len(text_dict)}"
                                text_dict[object_id] = smartart_text
                                
                                # Add to slide metadata
                                for slide_data in slide_metadata:
                                    if slide_data["slide_number"] == slide_num:
                                        slide_data["content"].append(smartart_text)
                                        break
                
                # Process WordArt and other graphics
                elif 'wordArt' in uri or 'diagram' in uri:
                    text = extract_text_from_element(graphic)
                    if text:
                        object_id = f"slide_{slide_num}_wordart_{len(text_dict)}"
                        text_dict[object_id] = text
                        
                        # Add to slide metadata
                        for slide_data in slide_metadata:
                            if slide_data["slide_number"] == slide_num:
                                slide_data["content"].append(text)
                                break
            
            # Find all text in the slide, including those that might be missed by python-pptx
            for para in slide_root.findall('.//a:p', namespaces):
                text = extract_text_from_element(para)
                # Only add text not already captured
                if text and not any(text in v for v in text_dict.values()):
                    object_id = f"slide_{slide_num}_text_{len(text_dict)}"
                    text_dict[object_id] = text
                    
                    # Add to slide metadata
                    for slide_data in slide_metadata:
                        if slide_data["slide_number"] == slide_num:
                            slide_data["content"].append(text)
                            break
    
    print(f"Deep extraction found {len(text_dict)} text elements")
    
    # Save to JSON if requested
    if output_json:
        output = {
            "text_dict": text_dict,
            "slide_metadata": slide_metadata
        }
        with open(output_json, 'w', encoding='utf-8') as f:
            json.dump(output, f, ensure_ascii=False, indent=2)
    
    return text_dict, slide_metadata

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Deep PowerPoint Text Extractor")
    parser.add_argument("--input-file", required=True, help="Input PowerPoint file (.pptx)")
    parser.add_argument("--output-json", help="Output JSON file for extracted text")
    
    args = parser.parse_args()
    
    deep_extract_text(args.input_file, args.output_json)