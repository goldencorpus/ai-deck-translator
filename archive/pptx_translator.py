import os
import json
import anthropic
from pptx import Presentation
import sys
import re
import time
import argparse
from datetime import datetime
from tqdm import tqdm
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

def extract_text(pptx_file):  
    """Extract text from PowerPoint presentation with enhanced support for various text content types"""
    prs = Presentation(pptx_file)  
    text_dict = {}
    slide_metadata = []  # Store structured context  
    
    # Extract text from a shape recursively, to handle grouped shapes
    def extract_shape_text(shape, shape_type="shape", parent_id=""):
        extracted_texts = {}
        
        # Get alt text if available (often contains important content)
        try:
            if hasattr(shape, "shape_properties") and hasattr(shape.shape_properties, "title") and shape.shape_properties.title:
                alt_text_id = f"{parent_id}_alt_text"
                extracted_texts[alt_text_id] = shape.shape_properties.title.strip()
        except:
            pass
        
        # Extract text from basic shape
        if hasattr(shape, "text") and shape.text.strip():
            shape_text = shape.text.strip()
            extracted_texts[parent_id] = shape_text
        
        # Handle SmartArt and other grouped shapes
        try:
            if hasattr(shape, "element") and hasattr(shape, "shapes"):
                # Extract text from any child shapes in a group or SmartArt
                try:
                    for i, child in enumerate(shape.shapes):
                        child_id = f"{parent_id}_child_{i}"
                        child_texts = extract_shape_text(child, "group_child", child_id)
                        extracted_texts.update(child_texts)
                except (AttributeError, TypeError, ValueError):
                    pass
        except:
            pass
        
        # Handle OLE objects (embedded Excel, etc.)
        try:
            if hasattr(shape, "shape_type") and shape.shape_type == 7:  # MSO_SHAPE_TYPE.OLE_OBJECT
                ole_id = f"{parent_id}_ole"
                # We can't directly extract text from OLE objects, but we can use alt text
                if hasattr(shape, "alternative_text") and shape.alternative_text:
                    extracted_texts[ole_id] = shape.alternative_text.strip()
        except:
            pass
            
        # Handle charts
        try:
            if hasattr(shape, "chart") and shape.chart:
                try:
                    chart_text = []
                    
                    # Extract chart title
                    if hasattr(shape.chart, "chart_title") and shape.chart.chart_title and shape.chart.chart_title.text_frame:
                        chart_text.append(shape.chart.chart_title.text_frame.text)
                    
                    # Extract category names
                    if hasattr(shape.chart, "plots"):
                        for plot in shape.chart.plots:
                            if hasattr(plot, "categories"):
                                for category in plot.categories:
                                    if category and category.strip():
                                        chart_text.append(category)
                    
                    if chart_text:
                        chart_id = f"{parent_id}_chart"
                        extracted_texts[chart_id] = " | ".join([t for t in chart_text if t])
                except (AttributeError, TypeError):
                    pass
        except:
            pass
        
        return extracted_texts
      
    for index, slide in enumerate(prs.slides):  
        slide_info = {  
            "slide_number": index + 1,  
            "title": "",  
            "content": []  
        }
        
        # Extract slide notes if any
        if hasattr(slide, "notes_slide") and slide.notes_slide:
            try:
                for note_shape in slide.notes_slide.shapes:
                    if hasattr(note_shape, "text") and note_shape.text.strip():
                        note_id = f"slide_{index+1}_notes"
                        note_text = note_shape.text.strip()
                        text_dict[note_id] = note_text
                        slide_info["content"].append(f"[Note: {note_text}]")
            except:
                pass
        
        # Process all shapes on the slide
        for shape_id, shape in enumerate(slide.shapes):
            base_id = f"slide_{index+1}_shape_{shape_id}"
            
            # Extract text from this shape (and any grouped sub-shapes)
            shape_texts = extract_shape_text(shape, "shape", base_id)
            
            # Add all extracted text to our dictionaries
            for obj_id, text_content in shape_texts.items():
                if text_content.strip():
                    text_dict[obj_id] = text_content.strip()
                    slide_info["content"].append(text_content.strip())
                    
                    # If this is the base shape and looks like a title
                    if obj_id == base_id:
                        # If this looks like a title shape, also add to slide title
                        if hasattr(shape, "is_title") and shape.is_title:
                            slide_info["title"] = text_content.strip()
                        # Alternatively check if it's a placeholder and has the right type
                        elif hasattr(shape, "is_placeholder") and shape.is_placeholder and hasattr(shape, "placeholder_format"):
                            try:
                                if shape.placeholder_format.type == 1:  # 1 is title placeholder
                                    slide_info["title"] = text_content.strip()
                            except (ValueError, AttributeError):
                                pass  # Handle case where placeholder_format raises an error
                        # Fallback for first shape on first slide
                        elif index == 0 and shape_id == 0:
                            slide_info["title"] = text_content.strip()
            
            # Handle tables separately
            if hasattr(shape, "has_table") and shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            # Create a unique ID for the table cell
                            cell_id = f"slide_{index+1}_table_{shape_id}_r{row_idx}_c{col_idx}"
                            text_dict[cell_id] = cell.text.strip()
                            slide_info["content"].append(cell.text.strip())
        
        # Try to extract header and footer text if available
        try:
            # Check placeholder shapes for header/footer content
            for shape in slide.shapes:
                if hasattr(shape, "is_placeholder") and shape.is_placeholder and hasattr(shape, "placeholder_format"):
                    try:
                        ph_type = shape.placeholder_format.type
                        # Common placeholder types for header (3), footer (4), date (2)
                        if ph_type in [2, 3, 4]:
                            if hasattr(shape, "text") and shape.text.strip():
                                ph_id = f"slide_{index+1}_placeholder_{ph_type}"
                                text_dict[ph_id] = shape.text.strip()
                                slide_info["content"].append(shape.text.strip())
                    except:
                        pass
        except:
            pass
            
        slide_metadata.append(slide_info)
    
    # Extract any text from master slides that might appear on all slides
    try:
        for idx, master in enumerate(prs.slide_masters):
            for shape_id, shape in enumerate(master.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    master_id = f"master_{idx+1}_shape_{shape_id}"
                    text_dict[master_id] = shape.text.strip()
    except:
        pass
        
    print(f"Enhanced extraction found {len(text_dict)} text elements")
    return text_dict, slide_metadata

def update_slides(pptx_file, output_file, translated_texts):
    """Update PowerPoint presentation with translated text, handling various content types"""
    prs = Presentation(pptx_file)
    updated_count = 0
    total_slides = len(prs.slides)
    
    # Function to update shape text recursively for grouped shapes
    def update_shape_text(shape, shape_type="shape", parent_id=""):
        shape_updated = False
        
        # Update basic shape text while preserving formatting
        if hasattr(shape, "text") and shape.text.strip():
            if parent_id in translated_texts:
                # If it has a text_frame, use that to preserve formatting
                if hasattr(shape, "text_frame"):
                    try:
                        # Save original text properties before replacement
                        para_properties = []
                        for para in shape.text_frame.paragraphs:
                            para_props = {
                                'font_size': None,
                                'font_bold': None,
                                'font_italic': None,
                                'font_underline': None,
                                'font_name': None,
                                'alignment': None,
                                'runs': []
                            }
                            
                            # Save paragraph level properties
                            if hasattr(para, "alignment") and para.alignment:
                                para_props['alignment'] = para.alignment
                                
                            # Save run level properties (character formatting)
                            for run in para.runs:
                                run_props = {}
                                if hasattr(run, "font") and run.font:
                                    if hasattr(run.font, "size") and run.font.size:
                                        run_props['size'] = run.font.size
                                    if hasattr(run.font, "bold") and run.font.bold is not None:
                                        run_props['bold'] = run.font.bold
                                    if hasattr(run.font, "italic") and run.font.italic is not None:
                                        run_props['italic'] = run.font.italic
                                    if hasattr(run.font, "underline") and run.font.underline is not None:
                                        run_props['underline'] = run.font.underline
                                    if hasattr(run.font, "name") and run.font.name:
                                        run_props['name'] = run.font.name
                                    if hasattr(run.font, "color") and run.font.color and hasattr(run.font.color, "rgb"):
                                        run_props['color'] = run.font.color.rgb
                                    
                                para_props['runs'].append(run_props)
                            
                            para_properties.append(para_props)
                        
                        # Clear the text frame and add the translated text
                        shape.text = ""
                        shape.text = translated_texts[parent_id]
                        
                        # Try to restore formatting if possible
                        if len(shape.text_frame.paragraphs) == len(para_properties):
                            for i, para in enumerate(shape.text_frame.paragraphs):
                                # Restore paragraph alignment
                                if para_properties[i]['alignment'] is not None:
                                    para.alignment = para_properties[i]['alignment']
                                
                                # Restore run level formatting if number of runs match
                                if len(para.runs) == len(para_properties[i]['runs']):
                                    for j, run in enumerate(para.runs):
                                        if hasattr(run, "font") and run.font:
                                            props = para_properties[i]['runs'][j]
                                            if 'size' in props and props['size']:
                                                run.font.size = props['size']
                                            if 'bold' in props and props['bold'] is not None:
                                                run.font.bold = props['bold']
                                            if 'italic' in props and props['italic'] is not None:
                                                run.font.italic = props['italic']
                                            if 'underline' in props and props['underline'] is not None:
                                                run.font.underline = props['underline']
                                            if 'name' in props and props['name']:
                                                run.font.name = props['name']
                                            if 'color' in props and props['color']:
                                                run.font.color.rgb = props['color']
                        shape_updated = True
                    except Exception as e:
                        # If detailed formatting preservation fails, fallback to simple replacement
                        shape.text = translated_texts[parent_id]
                        shape_updated = True
                else:
                    # Simple replacement if no text_frame
                    shape.text = translated_texts[parent_id]
                    shape_updated = True
        
        # Handle alt text
        alt_text_id = f"{parent_id}_alt_text"
        if alt_text_id in translated_texts and hasattr(shape, "shape_properties") and hasattr(shape.shape_properties, "title"):
            try:
                shape.shape_properties.title = translated_texts[alt_text_id]
                shape_updated = True
            except:
                pass
        
        # Handle SmartArt and other grouped shapes
        try:
            if hasattr(shape, "element") and hasattr(shape, "shapes"):
                # Update text from any child shapes in a group or SmartArt
                try:
                    for i, child in enumerate(shape.shapes):
                        child_id = f"{parent_id}_child_{i}"
                        child_updated = update_shape_text(child, "group_child", child_id)
                        shape_updated = shape_updated or child_updated
                except (AttributeError, TypeError, ValueError):
                    pass
        except:
            pass
        
        # Update OLE objects alt text if needed
        try:
            ole_id = f"{parent_id}_ole"
            if ole_id in translated_texts and hasattr(shape, "alternative_text"):
                shape.alternative_text = translated_texts[ole_id]
                shape_updated = True
        except:
            pass
            
        # Update chart text if possible
        try:
            chart_id = f"{parent_id}_chart"
            if chart_id in translated_texts and hasattr(shape, "chart") and shape.chart:
                try:
                    # We might not be able to update all chart elements, but we can try the title
                    if hasattr(shape.chart, "chart_title") and shape.chart.chart_title and shape.chart.chart_title.text_frame:
                        # For chart titles, take just the first part of the translated text
                        chart_title = translated_texts[chart_id].split(' | ')[0] if ' | ' in translated_texts[chart_id] else translated_texts[chart_id]
                        shape.chart.chart_title.text_frame.text = chart_title
                        shape_updated = True
                except:
                    pass
        except:
            pass
                
        return shape_updated
    
    # Try to update master slide text if it was extracted and translated
    try:
        for idx, master in enumerate(prs.slide_masters):
            for shape_id, shape in enumerate(master.shapes):
                master_id = f"master_{idx+1}_shape_{shape_id}"
                if master_id in translated_texts and hasattr(shape, "text") and shape.text.strip():
                    try:
                        shape.text = translated_texts[master_id]
                        updated_count += 1
                    except:
                        pass
    except:
        pass
    
    # Update each slide with progress bar
    from tqdm import tqdm
    with tqdm(total=total_slides, desc="Updating slides", unit="slide") as pbar:
        for index, slide in enumerate(prs.slides):
            # Update slide notes if any
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                try:
                    note_id = f"slide_{index+1}_notes"
                    if note_id in translated_texts:
                        for note_shape in slide.notes_slide.shapes:
                            if hasattr(note_shape, "text") and note_shape.text.strip():
                                note_shape.text = translated_texts[note_id]
                                updated_count += 1
            except:
                pass
                
            # Process all shapes on the slide
            for shape_id, shape in enumerate(slide.shapes):
                base_id = f"slide_{index+1}_shape_{shape_id}"
                
                # Update this shape (and any grouped sub-shapes)
                if update_shape_text(shape, "shape", base_id):
                    updated_count += 1
                
                # Handle tables separately
                if hasattr(shape, "has_table") and shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if cell.text.strip():
                                # Reconstruct cell ID to match the extraction phase
                                cell_id = f"slide_{index+1}_table_{shape_id}_r{row_idx}_c{col_idx}"
                                if cell_id in translated_texts:
                                    try:
                                        # Save original text properties before replacement
                                        para_properties = []
                                        for para in cell.text_frame.paragraphs:
                                        para_props = {
                                            'alignment': None,
                                            'runs': []
                                        }
                                        
                                        # Save paragraph level properties
                                        if hasattr(para, "alignment") and para.alignment:
                                            para_props['alignment'] = para.alignment
                                            
                                        # Save run level properties (character formatting)
                                        for run in para.runs:
                                            run_props = {}
                                            if hasattr(run, "font") and run.font:
                                                if hasattr(run.font, "size") and run.font.size:
                                                    run_props['size'] = run.font.size
                                                if hasattr(run.font, "bold") and run.font.bold is not None:
                                                    run_props['bold'] = run.font.bold
                                                if hasattr(run.font, "italic") and run.font.italic is not None:
                                                    run_props['italic'] = run.font.italic
                                                if hasattr(run.font, "underline") and run.font.underline is not None:
                                                    run_props['underline'] = run.font.underline
                                                if hasattr(run.font, "name") and run.font.name:
                                                    run_props['name'] = run.font.name
                                                if hasattr(run.font, "color") and run.font.color and hasattr(run.font.color, "rgb"):
                                                    run_props['color'] = run.font.color.rgb
                                            
                                            # Store hyperlink information if present
                                            if hasattr(run, "hyperlink") and run.hyperlink:
                                                run_props['hyperlink'] = run.hyperlink.address
                                                
                                            para_props['runs'].append(run_props)
                                        
                                        para_properties.append(para_props)
                                    
                                    # Clear the text and add the translated text
                                    cell.text = ""
                                    cell.text = translated_texts[cell_id]
                                    
                                    # Try to restore formatting if possible
                                    if len(cell.text_frame.paragraphs) == len(para_properties):
                                        for i, para in enumerate(cell.text_frame.paragraphs):
                                            # Restore paragraph alignment
                                            if para_properties[i]['alignment'] is not None:
                                                para.alignment = para_properties[i]['alignment']
                                            
                                            # Restore run level formatting if number of runs match
                                            if len(para.runs) == len(para_properties[i]['runs']):
                                                for j, run in enumerate(para.runs):
                                                    if hasattr(run, "font") and run.font:
                                                        props = para_properties[i]['runs'][j]
                                                        if 'size' in props and props['size']:
                                                            run.font.size = props['size']
                                                        if 'bold' in props and props['bold'] is not None:
                                                            run.font.bold = props['bold']
                                                        if 'italic' in props and props['italic'] is not None:
                                                            run.font.italic = props['italic']
                                                        if 'underline' in props and props['underline'] is not None:
                                                            run.font.underline = props['underline']
                                                        if 'name' in props and props['name']:
                                                            run.font.name = props['name']
                                                        if 'color' in props and props['color']:
                                                            run.font.color.rgb = props['color']
                                                            
                                                        # Restore hyperlink if present
                                                        if 'hyperlink' in props and props['hyperlink'] and hasattr(run, "hyperlink"):
                                                            run.hyperlink.address = props['hyperlink']
                                except Exception as e:
                                    # Fall back to simple replacement if formatting preservation fails
                                    cell.text = translated_texts[cell_id]
                                
                                updated_count += 1
        
            # Try to update header and footer text if available
            try:
                # Check placeholder shapes for header/footer content
                for shape in slide.shapes:
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder and hasattr(shape, "placeholder_format"):
                        try:
                            ph_type = shape.placeholder_format.type
                            if ph_type in [2, 3, 4]:  # header, footer, date placeholders
                                ph_id = f"slide_{index+1}_placeholder_{ph_type}"
                                if ph_id in translated_texts and hasattr(shape, "text") and shape.text.strip():
                                    shape.text = translated_texts[ph_id]
                                    updated_count += 1
                    except:
                        pass
            except:
                pass
                
            # Update progress bar after each slide
            completion_percentage = int(100 * (index + 1) / total_slides)
            pbar.set_description(f"Updating slides: {completion_percentage}% complete")
            pbar.update(1)
            pbar.refresh()
    
    print(f"Updated {updated_count} text elements in the presentation")
    # Save the updated presentation to a new file
    prs.save(output_file)
    return output_file

def split_dict_into_smart_batches(input_dict, max_input_tokens=150000, prompt_tokens=2000):
    """
    Split a dictionary into batches based on estimated token count to optimize API usage.
    """
    # Function to estimate tokens in a string with better Unicode/multibyte handling
    def estimate_tokens(text):
        if text is None:
            return 0
            
        # Convert to string if not already
        text_str = str(text)
        
        # For Asian languages (CJK), use 1.5 characters per token as a conservative estimate
        cjk_chars = sum(1 for c in text_str if ord(c) > 0x4E00 and ord(c) < 0x9FFF)  # Basic CJK Unified Ideographs
        ascii_chars = len(text_str) - cjk_chars
        
        # 4 ASCII chars per token, ~1.5 CJK chars per token (rough estimate)
        token_estimate = (ascii_chars // 4) + (cjk_chars // 1.5) + 1  # Add 1 to round up
        
        return int(token_estimate)
    
    items = list(input_dict.items())
    batches = []
    current_batch = {}
    current_token_count = prompt_tokens
    
    # Sort items by estimated token length (optional)
    items.sort(key=lambda x: estimate_tokens(x[1]), reverse=True)
    
    for key, value in items:
        item_tokens = estimate_tokens(key) + estimate_tokens(value) + 10  # +10 for JSON formatting
        
        if current_token_count + item_tokens > max_input_tokens and current_batch:
            batches.append(current_batch)
            current_batch = {}
            current_token_count = prompt_tokens
        
        current_batch[key] = value
        current_token_count += item_tokens
    
    if current_batch:
        batches.append(current_batch)
    
    total_items = len(input_dict)
    batch_sizes = [len(batch) for batch in batches]
    avg_batch_size = sum(batch_sizes) / len(batches) if batches else 0
    
    print(f"Created {len(batches)} batches from {total_items} items")
    print(f"Batch sizes: min={min(batch_sizes) if batches else 0}, max={max(batch_sizes) if batches else 0}, avg={avg_batch_size:.1f}")
    print(f"Estimated token usage efficiency: {(sum(batch_sizes)/total_items)*100:.1f}%")
    
    return batches

def repair_json(json_content):
    """
    More robust JSON repair function that can handle various common issues including Unicode.
    """
    original_content = json_content
    try:
        return json.loads(json_content)
    except json.JSONDecodeError as e:
        print(f"Initial JSON parsing error: {e}")
        
        # Try to handle Unicode/multibyte character issues
        try:
            # Ensure the content is properly encoded
            if isinstance(json_content, bytes):
                json_content = json_content.decode('utf-8', errors='replace')
            
            # Remove or replace invisible/control characters that might cause issues
            json_content = re.sub(r'[\x00-\x1F\x7F]', '', json_content)
        except Exception as enc_err:
            print(f"Error handling encoding: {enc_err}")
        
        if "Unterminated string" in str(e):
            error_info = str(e)
            line_match = re.search(r'line (\d+)', error_info)
            col_match = re.search(r'column (\d+)', error_info)
            
            if line_match and col_match:
                line_num = int(line_match.group(1))
                col_num = int(col_match.group(1))
                lines = json_content.split('\n')
                if 0 <= line_num-1 < len(lines):
                    line = lines[line_num-1]
                    if col_num-1 < len(line) and line[col_num-1] == '"':
                        line = line[:col_num-1] + '\\"' + line[col_num:]
                    else:
                        line = line + '"'
                    lines[line_num-1] = line
                    json_content = '\n'.join(lines)
        
        # Handle improperly escaped characters and unicode escapes
        json_content = re.sub(r'(?<!\\)\\(?!["\\/bfnrt]|u[0-9a-fA-F]{4})', r'\\\\', json_content)
        
        # Fix unbalanced braces
        brace_count = json_content.count('{') - json_content.count('}')
        if brace_count > 0:
            json_content = json_content + ('}' * brace_count)
        elif brace_count < 0:
            for _ in range(-brace_count):
                json_content = json_content.rstrip().rstrip('}').rstrip()
        
        # Fix common JSON syntax issues
        json_content = re.sub(r',\s*}', '}', json_content)
        json_content = re.sub(r',\s*]', ']', json_content)
        
        # Ensure property names are properly quoted
        def fix_property_names(match):
            prop = match.group(1)
            if not (prop.startswith('"') and prop.endswith('"')):
                return f'"{prop}":'
            return match.group(0)
        
        json_content = re.sub(r'([a-zA-Z0-9_]+):', fix_property_names, json_content)
        
        # Try to parse the repaired JSON
        try:
            return json.loads(json_content)
        except json.JSONDecodeError as e2:
            print(f"JSON repair attempt failed: {e2}")
            
            # Fallback: extract key-value pairs using regex
            result = {}
            # Modified pattern to handle Unicode characters better
            pattern = r'"((?:[^"\\]|\\.)*)"\s*:\s*"((?:[^"\\]|\\.)*)"'
            for match in re.finditer(pattern, original_content):
                try:
                    key, value = match.groups()
                    # Unescape escaped quotes in the extracted strings
                    key = key.replace('\\"', '"')
                    value = value.replace('\\"', '"')
                    result[key] = value
                except Exception as extract_err:
                    print(f"Error extracting key-value pair: {extract_err}")
            
            # Also try to capture numeric values
            num_pattern = r'"((?:[^"\\]|\\.)*)"\s*:\s*([0-9]+(?:\.[0-9]+)?)'
            for match in re.finditer(num_pattern, original_content):
                try:
                    key, value = match.groups()
                    key = key.replace('\\"', '"')
                    # Convert to int or float as appropriate
                    if '.' in value:
                        result[key] = float(value)
                    else:
                        result[key] = int(value)
                except Exception as extract_err:
                    print(f"Error extracting numeric value: {extract_err}")
            
            if result:
                print(f"Managed to extract {len(result)} key-value pairs through regex")
                return result
            
            raise e

def extract_json_blocks(text):
    """
    Extract valid JSON blocks from text that might contain multiple partial JSON objects.
    Enhanced to better handle Unicode and Japanese characters.
    """
    # First try to find complete JSON objects
    try:
        # Clean the text: remove any leading/trailing non-JSON content
        text = re.sub(r'^[^{]*', '', text)  # Remove anything before the first {
        text = re.sub(r'[^}]*$', '', text)  # Remove anything after the last }
        
        # Try extracting json blocks with a more robust pattern
        # This pattern tries to match balanced { } pairs
        depth = 0
        start = -1
        blocks = []
        
        for i, char in enumerate(text):
            if char == '{':
                if depth == 0:
                    start = i
                depth += 1
            elif char == '}':
                depth -= 1
                if depth == 0 and start != -1:
                    blocks.append(text[start:i+1])
                    start = -1
        
        if not blocks:
            # Fallback to simpler regex if the balanced matching didn't work
            potential_blocks = re.findall(r'({[^{]*?})', text)
            blocks = potential_blocks
        
        valid_blocks = []
        for block in blocks:
            try:
                # Fix common issues that might occur in the JSON block
                block = re.sub(r',\s*}', '}', block)  # Remove trailing commas
                block = re.sub(r'([a-zA-Z0-9_]+):', r'"\1":', block)  # Quote unquoted keys
                
                parsed = json.loads(block)
                valid_blocks.append(parsed)
            except json.JSONDecodeError:
                try:
                    # Try to repair the JSON before giving up
                    repaired = repair_json(block)
                    if repaired:
                        valid_blocks.append(repaired)
                except:
                    pass
        
        if valid_blocks:
            combined = {}
            for block in valid_blocks:
                combined.update(block)
            return combined
            
    except Exception as e:
        print(f"Error in JSON block extraction: {e}")
        
    # Fallback: just try a direct key-value extraction
    try:
        result = {}
        # Look for key-value pairs directly, handling Unicode properly
        pattern = r'"((?:[^"\\]|\\.)*)"\s*:\s*"((?:[^"\\]|\\.)*)"'
        for match in re.finditer(pattern, text):
            key, value = match.groups()
            # Unescape escaped quotes
            key = key.replace('\\"', '"')
            value = value.replace('\\"', '"')
            result[key] = value
        
        if result:
            print(f"Direct key-value extraction found {len(result)} pairs")
            return result
    except Exception as e:
        print(f"Fallback extraction error: {e}")
        
    return None

def setup_recovery_system(file_id, text_dict, slide_metadata, source_language, target_language, resume_file=None):
    """
    Set up a recovery system for batch processing.
    """
    recovery_dir = "translation_recovery"
    os.makedirs(recovery_dir, exist_ok=True)
    
    if resume_file and os.path.exists(resume_file):
        with open(resume_file, 'r', encoding='utf-8') as f:
            recovery_state = json.load(f)
        print(f"Resuming translation from recovery file: {resume_file}")
        recovery_file = resume_file
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        recovery_file = os.path.join(recovery_dir, f"recovery_pptx_{file_id}_{timestamp}.json")
        recovery_state = {
            "file_id": file_id,
            "completed_batches": [],
            "failed_batches": [],
            "translated_items": {},
            "source_language": source_language,
            "target_language": target_language,
            "total_items": len(text_dict),
            "start_time": timestamp,
            "last_updated": timestamp
        }
        with open(recovery_file, 'w', encoding='utf-8') as f:
            json.dump(recovery_state, f, ensure_ascii=False, indent=2)
        print(f"Created new recovery file: {recovery_file}")
    
    def save_recovery_state():
        recovery_state["last_updated"] = datetime.now().strftime("%Y%m%d_%H%M%S")
        with open(recovery_file, 'w', encoding='utf-8') as f:
            json.dump(recovery_state, f, ensure_ascii=False, indent=2)
    
    return recovery_state, recovery_file, save_recovery_state

def translate_batch(batch, batch_index, slide_metadata, source_language, target_language, api_key=None, max_retries=2, cost_tracker=None):
    """
    Translate a single batch with retry logic.
    """
    batch_copy = batch.copy()
    
    def clean_text(text):
        return text.replace('\\n', '\n').replace('\\u000b', '\v').replace('\\t', '\t')
        
    # Calculate estimated cost based on prompt and completion tokens
    def estimate_cost(prompt_tokens, completion_tokens, model="claude-3-7-sonnet"):
        # Claude 3.5 Sonnet pricing: $3 per 1M input tokens, $15 per 1M output tokens
        # Claude 3 Sonnet pricing: $3 per 1M input tokens, $15 per 1M output tokens
        rates = {
            "claude-3-7-sonnet": {"input": 3.0 / 1000000, "output": 15.0 / 1000000},
            "claude-3-sonnet": {"input": 3.0 / 1000000, "output": 15.0 / 1000000}, 
            "claude-3-opus": {"input": 15.0 / 1000000, "output": 75.0 / 1000000},
            "claude-3-haiku": {"input": 0.25 / 1000000, "output": 1.25 / 1000000}
        }
        
        model_rates = rates.get(model, rates["claude-3-7-sonnet"])  # Default to sonnet if not found
        input_cost = prompt_tokens * model_rates["input"]
        output_cost = completion_tokens * model_rates["output"]
        total_cost = input_cost + output_cost
        
        return {
            "input_tokens": prompt_tokens,
            "output_tokens": completion_tokens,
            "input_cost": input_cost,
            "output_cost": output_cost,
            "total_cost": total_cost
        }
    
    # Use provided API key or fall back to environment variable
    api_key = api_key or os.getenv("CLAUDE_API_KEY")
    if not api_key:
        raise ValueError("Claude API key must be provided either as an argument or via CLAUDE_API_KEY environment variable")
        
    client = anthropic.Anthropic(
        api_key=api_key,
        default_headers={
            "anthropic-beta": "output-128k-2025-02-19"
        }
    )
    
    structured_context = json.dumps(slide_metadata, ensure_ascii=False, indent=2)
    
    system_prompt = f"""You are a professional translator. Translate from {source_language} to {target_language}.
Ensure consistency in terminology and contextual meaning.

IMPORTANT: If you encounter text that appears to already be in {target_language}, preserve it exactly as is without any changes.
Do not translate text that is already in {target_language}.

PRIVACY NOTICE: Do not store, learn from, or retain any of the content provided for translation.
This is confidential material that should only be processed for immediate translation purposes."""
    
    user_message = f"""
Translate the following JSON object from {source_language} to {target_language}.
Consider the structured slide context provided below for context and consistency.

IMPORTANT INSTRUCTIONS:
- If any text appears to already be in {target_language}, keep it exactly as is.
- Only translate text that is in {source_language}.
- Do NOT include escape sequences for newlines (\\n) or other characters - use the actual characters.
- Return VALID JSON format with all keys and values properly enclosed in double quotes.
- Ensure all property names and string values are properly quoted with double quotes.
- Do not include any trailing commas.

This is batch {batch_index} with {len(batch_copy)} items.

Slide Context:
{structured_context}

Now translate the following structured JSON object while preserving its format:
{json.dumps(batch_copy, ensure_ascii=False, indent=2)}

Reply ONLY with the translated JSON. The JSON MUST be valid and parseable.
"""

    for retry in range(max_retries + 1):
        try:
            response = client.messages.create(
                model="claude-3-7-sonnet-20250219",
                system=system_prompt,
                max_tokens=4000,
                messages=[
                    {"role": "user", "content": user_message}
                ],
                metadata={
                    "user_id": "anonymous_user"
                }
            )
            
            # Track token usage and cost
            prompt_tokens = response.usage.input_tokens
            completion_tokens = response.usage.output_tokens
            
            batch_cost = estimate_cost(prompt_tokens, completion_tokens, "claude-3-7-sonnet")
            
            # Update the cost tracker if provided
            if cost_tracker is not None:
                cost_tracker["total_input_tokens"] += prompt_tokens
                cost_tracker["total_output_tokens"] += completion_tokens
                cost_tracker["total_input_cost"] += batch_cost["input_cost"]
                cost_tracker["total_output_cost"] += batch_cost["output_cost"]
                cost_tracker["total_cost"] += batch_cost["total_cost"]
                cost_tracker["api_calls"] += 1
            
            print(f"Batch {batch_index} token usage: {prompt_tokens} input + {completion_tokens} output tokens")
            print(f"Batch {batch_index} cost: ${batch_cost['total_cost']:.4f} (${batch_cost['input_cost']:.4f} input + ${batch_cost['output_cost']:.4f} output)")
            
            # Show running total cost
            if cost_tracker is not None:
                print(f"Running total: ${cost_tracker['total_cost']:.4f} for {cost_tracker['api_calls']} API calls")

            translated_text = response.content[0].text
            
            if "```json" in translated_text:
                json_content = translated_text.split("```json")[1].split("```")[0].strip()
            elif "```" in translated_text:
                json_content = translated_text.split("```")[1].strip()
            else:
                json_content = translated_text.strip()
            
            try:
                batch_result = json.loads(json_content)
            except json.JSONDecodeError as e:
                try:
                    batch_result = repair_json(json_content)
                except Exception as e2:
                    extracted_result = extract_json_blocks(json_content)
                    if extracted_result:
                        batch_result = extracted_result
                        print(f"Extracted {len(batch_result)} items through JSON block extraction")
                    else:
                        if retry < max_retries:
                            print(f"JSON parsing failed on attempt {retry+1}, retrying...")
                            time.sleep(3)
                            continue
                        else:
                            raise e
            
            for key, value in batch_result.items():
                if isinstance(value, str):
                    batch_result[key] = clean_text(value)
            
            print(f"Successfully processed batch {batch_index}")
            return batch_result
                
        except Exception as e:
            if retry < max_retries:
                print(f"Error in batch {batch_index} (attempt {retry+1}): {e}")
                print(f"Retrying in 5 seconds...")
                time.sleep(5)
            else:
                print(f"All {max_retries + 1} attempts failed for batch {batch_index}: {e}")
                raise e

def translate_text(text_dict, slide_metadata, source_language, target_language, resume_file=None, api_key=None):
    # Use provided API key or fall back to environment variable
    api_key = api_key or os.getenv("CLAUDE_API_KEY")
    if not api_key:
        raise ValueError("Claude API key must be provided either as an argument or via CLAUDE_API_KEY environment variable")
        
    client = anthropic.Anthropic(
        api_key=api_key,
        default_headers={
            "anthropic-beta": "output-128k-2025-02-19"
        }
    )
    
    # Initialize cost tracking
    cost_tracker = {
        "total_input_tokens": 0,
        "total_output_tokens": 0,
        "total_input_cost": 0.0,
        "total_output_cost": 0.0,
        "total_cost": 0.0,
        "api_calls": 0
    }
    
    def deduplicate_content(input_dict):
        content_to_keys = {}
        for key, value in input_dict.items():
            if value in content_to_keys:
                content_to_keys[value].append(key)
            else:
                content_to_keys[value] = [key]
        
        unique_content = {}
        duplicates_map = {}
        
        for content, keys in content_to_keys.items():
            representative_key = keys[0]
            unique_content[representative_key] = content
            for key in keys:
                duplicates_map[key] = representative_key
        
        print(f"Found {len(input_dict) - len(unique_content)} duplicate content items")
        print(f"Reduced from {len(input_dict)} to {len(unique_content)} unique content items to translate")
        
        return unique_content, duplicates_map
    
    # Use a simple file ID based on timestamp if not provided
    file_id = os.path.basename(str(time.time()).replace('.', ''))
    
    recovery_state, recovery_file, save_recovery_state = setup_recovery_system(
        file_id, text_dict, slide_metadata, source_language, target_language, resume_file
    )
    
    if not recovery_state["translated_items"]:
        unique_text_dict, duplicates_map = deduplicate_content(text_dict)
        recovery_state["duplicates_map"] = duplicates_map
        save_recovery_state()
    else:
        unique_text_dict = {k: text_dict[k] for k in text_dict.keys() 
                          if k not in [rep_key for rep_key in recovery_state["duplicates_map"].values()]}
        duplicates_map = recovery_state["duplicates_map"]
        print(f"Resumed with {len(recovery_state['translated_items'])} already translated items")
    
    remaining_dict = {k: v for k, v in unique_text_dict.items() 
                     if k not in recovery_state["translated_items"]}
    
    if not remaining_dict:
        print("All items have already been translated. Nothing to do.")
        full_translated_dict = recovery_state["translated_items"].copy()
    else:
        # For Japanese or other multibyte languages, use smaller batches
        max_tokens = 80000 if target_language in ["ja", "zh", "ko"] else 150000
        prompt_tokens = 2000
        
        print(f"Using smaller batch size for {target_language} translation" if target_language in ["ja", "zh", "ko"] else "Using standard batch size")
        batches = split_dict_into_smart_batches(remaining_dict, max_input_tokens=max_tokens, prompt_tokens=prompt_tokens)
        print(f"Splitting translation into {len(batches)} batches")
        
        unique_translated_dict = recovery_state["translated_items"].copy()
        
        with tqdm(total=len(batches), desc="Translating", unit="batch") as pbar:
            for batch_index, batch in enumerate(batches):
                batch_id = f"batch_{batch_index+1}"
                if batch_id in recovery_state["completed_batches"]:
                    print(f"Skipping already completed batch {batch_id}")
                    pbar.update(1)
                    continue
                
                print(f"\nProcessing batch {batch_index+1} of {len(batches)} with {len(batch)} items...")
                
                try:
                    batch_result = translate_batch(
                        batch, batch_index+1, slide_metadata, 
                        source_language, target_language, api_key=api_key,
                        cost_tracker=cost_tracker
                    )
                    
                    unique_translated_dict.update(batch_result)
                    recovery_state["translated_items"].update(batch_result)
                    recovery_state["completed_batches"].append(batch_id)
                    save_recovery_state()
                    
                except Exception as e:
                    print(f"Error in batch {batch_index+1}: {e}")
                    recovery_state["failed_batches"].append({
                        "batch_id": batch_id,
                        "keys": list(batch.keys()),
                        "error": str(e)
                    })
                    save_recovery_state()
                    print("Continuing with next batch...")
                
                pbar.update(1)
                completion_percentage = int(100 * (batch_index + 1) / len(batches))
                pbar.set_description(f"Translating: {completion_percentage}% complete")
                # Force refresh the progress bar display
                pbar.refresh()
        
        print(f"\nTranslation of unique content completed with {len(unique_translated_dict)} items out of {len(unique_text_dict)} unique items")
        
        if recovery_state["failed_batches"]:
            print(f"\nRetrying {len(recovery_state['failed_batches'])} failed batches with smaller chunks...")
            
            for failed_batch in list(recovery_state["failed_batches"]):
                batch_id = failed_batch["batch_id"]
                keys = failed_batch["keys"]
                
                retry_batch = {k: text_dict[k] for k in keys if k in text_dict}
                # Use smaller chunks for CJK languages
                divisor = 8 if target_language in ["ja", "zh", "ko"] else 4
                chunk_size = max(5, len(retry_batch) // divisor)
                retry_items = list(retry_batch.items())
                sub_batches = [dict(retry_items[i:i+chunk_size]) 
                               for i in range(0, len(retry_items), chunk_size)]
                
                print(f"Split failed batch {batch_id} into {len(sub_batches)} smaller chunks of size ~{chunk_size}")
                
                for i, sub_batch in enumerate(sub_batches):
                    sub_id = f"{batch_id}_sub_{i+1}"
                    
                    try:
                        print(f"Processing sub-batch {i+1}/{len(sub_batches)} for failed batch {batch_id}")
                        sub_result = translate_batch(
                            sub_batch, f"{batch_id}.{i+1}", slide_metadata, 
                            source_language, target_language, api_key=api_key, 
                            max_retries=3, cost_tracker=cost_tracker
                        )
                        unique_translated_dict.update(sub_result)
                        recovery_state["translated_items"].update(sub_result)
                        recovery_state["completed_batches"].append(sub_id)
                        save_recovery_state()
                        
                    except Exception as e:
                        print(f"Error in sub-batch {i+1} of failed batch {batch_id}: {e}")
                        continue
                
                recovery_state["failed_batches"].remove(failed_batch)
                save_recovery_state()
        
        full_translated_dict = {}
        for key, value in unique_translated_dict.items():
            full_translated_dict[key] = value
        
        for original_key, rep_key in duplicates_map.items():
            if rep_key in unique_translated_dict and original_key != rep_key:
                full_translated_dict[original_key] = unique_translated_dict[rep_key]
    
    print(f"Reconstructed full translation dictionary with {len(full_translated_dict)} items")
    
    missing_keys = set(text_dict.keys()) - set(full_translated_dict.keys())
    if missing_keys:
        print(f"Warning: {len(missing_keys)} keys were not translated: {list(missing_keys)[:5]}...")
        if len(missing_keys) > 0:
            print(f"Attempting to translate {len(missing_keys)} missing keys in a final batch...")
            missing_dict = {k: text_dict[k] for k in missing_keys if k in text_dict}
            
            try:
                structured_context = json.dumps(slide_metadata, ensure_ascii=False, indent=2)
                
                system_prompt = f"""You are a professional translator. Translate from {source_language} to {target_language}.
Ensure consistency in terminology and contextual meaning.

IMPORTANT: If you encounter text that appears to already be in {target_language}, preserve it exactly as is.
Do not translate text that is already in {target_language}."""
                
                user_message = f"""
Translate the following JSON object from {source_language} to {target_language}.
This is a final batch to catch any missing translations.

IMPORTANT INSTRUCTIONS:
- If any text appears to already be in {target_language}, keep it exactly as is.
- Only translate text that is in {source_language}.
- Do NOT include escape sequences for newlines (\\n) or other characters - use the actual characters.
- Return VALID JSON format with all keys and values properly enclosed in double quotes.

Now translate the following structured JSON object:
{json.dumps(missing_dict, ensure_ascii=False, indent=2)}

Reply ONLY with the translated JSON.
"""
                
                response = client.messages.create(
                    model="claude-3-7-sonnet-20250219",
                    system=system_prompt,
                    max_tokens=4000,
                    messages=[
                        {"role": "user", "content": user_message}
                    ],
                    metadata={
                        "user_id": "anonymous_user"
                    }
                )
                
                # Track token usage and cost for final batch
                prompt_tokens = response.usage.input_tokens
                completion_tokens = response.usage.output_tokens
                
                # Calculate cost
                final_cost = {
                    "input_cost": prompt_tokens * (3.0 / 1000000),
                    "output_cost": completion_tokens * (15.0 / 1000000),
                    "total_cost": (prompt_tokens * (3.0 / 1000000)) + (completion_tokens * (15.0 / 1000000))
                }
                
                # Update cost tracker
                cost_tracker["total_input_tokens"] += prompt_tokens
                cost_tracker["total_output_tokens"] += completion_tokens
                cost_tracker["total_input_cost"] += final_cost["input_cost"]
                cost_tracker["total_output_cost"] += final_cost["output_cost"]
                cost_tracker["total_cost"] += final_cost["total_cost"]
                cost_tracker["api_calls"] += 1
                
                print(f"Final batch token usage: {prompt_tokens} input + {completion_tokens} output tokens")
                print(f"Final batch cost: ${final_cost['total_cost']:.4f}")
                print(f"Running total: ${cost_tracker['total_cost']:.4f} for {cost_tracker['api_calls']} API calls")
                
                translated_text = response.content[0].text
                
                if "```json" in translated_text:
                    json_content = translated_text.split("```json")[1].split("```")[0].strip()
                elif "```" in translated_text:
                    json_content = translated_text.split("```")[1].strip()
                else:
                    json_content = translated_text.strip()
                
                try:
                    final_batch = json.loads(json_content)
                except json.JSONDecodeError:
                    try:
                        final_batch = repair_json(json_content)
                    except:
                        extracted = extract_json_blocks(json_content)
                        if extracted:
                            final_batch = extracted
                        else:
                            raise
                
                def clean_text(text):
                    return text.replace('\\n', '\n').replace('\\u000b', '\v').replace('\\t', '\t')
                
                for key, value in final_batch.items():
                    if isinstance(value, str):
                        final_batch[key] = clean_text(value)
                
                full_translated_dict.update(final_batch)
                recovery_state["translated_items"].update(final_batch)
                save_recovery_state()
                
                print(f"Successfully processed final batch with {len(final_batch)} additional items")
                
                missing_keys = set(text_dict.keys()) - set(full_translated_dict.keys())
                if missing_keys:
                    print(f"Final warning: {len(missing_keys)} keys still not translated: {list(missing_keys)[:5]}...")
                else:
                    print("All items successfully translated!")
                    
            except Exception as e:
                print(f"Failed to process final batch: {e}")
    else:
        print("All items successfully translated!")
    
    # Print final cost summary
    print("\n=== API Cost Summary ===")
    print(f"Total API calls: {cost_tracker['api_calls']}")
    print(f"Total input tokens: {cost_tracker['total_input_tokens']:,}")
    print(f"Total output tokens: {cost_tracker['total_output_tokens']:,}")
    print(f"Total tokens: {cost_tracker['total_input_tokens'] + cost_tracker['total_output_tokens']:,}")
    print(f"Input cost: ${cost_tracker['total_input_cost']:.4f}")
    print(f"Output cost: ${cost_tracker['total_output_cost']:.4f}")
    print(f"Total cost: ${cost_tracker['total_cost']:.4f}")
    
    return full_translated_dict

def translate_pptx(input_file, output_file, source_language="en", target_language="fr", resume_file=None, api_key=None):
    """Main function to translate PowerPoint files"""
    print(f"Extracting text from {input_file}...")
    text_dict, slide_metadata = extract_text(input_file)
    print(f"Found {len(text_dict)} text elements across {len(slide_metadata)} slides")
    
    print(f"Translating from {source_language} to {target_language}...")
    translated_texts = translate_text(text_dict, slide_metadata, source_language, target_language, resume_file, api_key=api_key)
    
    print(f"Updating PowerPoint with translated text...")
    update_slides(input_file, output_file, translated_texts)
    
    print(f"Translation completed!")
    print(f"Translated presentation saved as: {output_file}")
    
    return output_file

def list_recovery_files():
    """List all available recovery files and their status"""
    recovery_dir = "translation_recovery"
    if not os.path.exists(recovery_dir):
        print("No recovery directory found.")
        return
    
    recovery_files = [f for f in os.listdir(recovery_dir) if f.startswith("recovery_pptx_") and f.endswith(".json")]
    
    if not recovery_files:
        print("No recovery files found.")
        return
    
    print(f"Found {len(recovery_files)} recovery files:")
    for f in recovery_files:
        file_path = os.path.join(recovery_dir, f)
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                total = data.get("total_items", 0)
                translated = len(data.get("translated_items", {}))
                failed = len(data.get("failed_batches", []))
                progress = (translated / total * 100) if total > 0 else 0
                
                print(f"  {f}")
                print(f"    Progress: {progress:.1f}% ({translated}/{total} items)")
                print(f"    Failed batches: {failed}")
                print(f"    Start time: {data.get('start_time', 'unknown')}")
                print(f"    Last updated: {data.get('last_updated', 'unknown')}")
                print()
        except Exception as e:
            print(f"  {f} - Error reading file: {e}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PowerPoint PPTX Translator")
    parser.add_argument("--resume", help="Resume translation from a recovery file")
    parser.add_argument("--list-recovery", action="store_true", help="List available recovery files")
    parser.add_argument("--input-file", help="Input PowerPoint file (.pptx)")
    parser.add_argument("--output-file", help="Output PowerPoint file (.pptx)")
    parser.add_argument("--source-language", help="Source language code (e.g., en)")
    parser.add_argument("--target-language", help="Target language code (e.g., ja)")
    parser.add_argument("--api-key", help="Claude API Key (can also be set as CLAUDE_API_KEY environment variable)")
    
    args = parser.parse_args()
    
    if args.list_recovery:
        list_recovery_files()
        sys.exit(0)
    
    if args.input_file:
        input_file = args.input_file
    else:
        input_file = input("Enter path to PowerPoint file (.pptx): ")
    
    if not os.path.exists(input_file):
        print(f"Error: Input file '{input_file}' does not exist")
        sys.exit(1)
    
    if args.output_file:
        output_file = args.output_file
    else:
        # Create a default output filename based on the input filename
        base_name = os.path.splitext(os.path.basename(input_file))[0]
        if args.target_language:
            output_file = f"{base_name}_{args.target_language}.pptx"
        else:
            output_file = f"{base_name}_translated.pptx"
    
    if args.source_language:
        source_language = args.source_language
    else:
        source_language = input("Enter source language (e.g., en for English): ")
    
    if args.target_language:
        target_language = args.target_language
    else:
        target_language = input("Enter target language (e.g., fr for French): ")
    
    translate_pptx(input_file, output_file, source_language, target_language, args.resume, api_key=args.api_key)