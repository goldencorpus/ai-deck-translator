import sys
from pptx import Presentation
import json

try:
    # Load the presentation
    prs = Presentation(sys.argv[1])
    
    # Extract text from each slide
    slides_text = {}
    slide_metadata = {"slides": []}
    slide_index = 0
    
    for slide in prs.slides:
        slide_index += 1
        slide_key = f"slide{slide_index}"
        slides_text[slide_key] = {}
        
        # Create slide metadata
        slide_meta = {
            "id": slide_key,
            "elements": []
        }
        
        # Get text from shapes
        shape_index = 0
        for shape in slide.shapes:
            shape_index += 1
            shape_id = f"{slide_key}_shape{shape_index}"
            
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                print(f"Found table in slide {slide_index}, shape {shape_index}")
                
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            # Create unique ID for table cell
                            cell_id = f"{shape_id}_table_r{row_idx}c{col_idx}"
                            slides_text[slide_key][cell_id] = cell.text
                            
                            # Store metadata for table cell
                            element_meta = {
                                "id": cell_id,
                                "type": "table_cell",
                                "row": row_idx,
                                "column": col_idx,
                                "parent_shape": shape_id
                            }
                            slide_meta["elements"].append(element_meta)
            
            # Handle regular text
            elif hasattr(shape, "text") and shape.text:
                # Try to identify the shape type
                shape_type = "text"
                if shape.name.lower().startswith("title"):
                    shape_type = "title"
                elif slide_index == 1 and shape_index == 2:
                    shape_type = "subtitle"
                
                text_id = f"{shape_id}_{shape_type}"
                slides_text[slide_key][text_id] = shape.text
                
                # Store metadata for text
                element_meta = {
                    "id": text_id,
                    "type": shape_type,
                    "parent_shape": shape_id
                }
                slide_meta["elements"].append(element_meta)
        
        # Add this slide's metadata to the overall metadata
        slide_metadata["slides"].append(slide_meta)
    
    # Create result with both text and metadata
    result = {
        "text_elements": slides_text,
        "metadata": slide_metadata
    }
    
    # Print as JSON
    print(json.dumps(result, indent=2))
    print("Extraction successful!")
    
except Exception as e:
    print(f"Error during extraction: {str(e)}")
    sys.exit(1)