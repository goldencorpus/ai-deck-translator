import sys
from pptx import Presentation
import json

try:
    # Load the presentation and the translations
    prs = Presentation(sys.argv[1])
    
    with open(sys.argv[2], 'r') as f:
        translation_data = json.load(f)
    
    translations = translation_data['translations']
    
    print("Slide texts found in the presentation:")
    
    # First, display what we found in the presentation
    slide_index = 0
    for slide in prs.slides:
        slide_index += 1
        print(f"Slide {slide_index}:")
        
        # Track shape index for this slide
        shape_index = 0
        
        for shape in slide.shapes:
            shape_index += 1
            shape_id = f"slide{slide_index}_shape{shape_index}"
            
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                print(f"  - Table: {shape.name}")
                
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            cell_id = f"{shape_id}_table_r{row_idx}c{col_idx}"
                            print(f"    - Cell [{row_idx},{col_idx}]: '{cell.text}'")
            
            # Handle regular text
            elif hasattr(shape, "text") and shape.text:
                # Replace newlines for display only
                display_text = shape.text.replace("\n", " ")
                print(f"  - {shape.name}: '{display_text}'")
    
    # Reset and apply translations
    print("\nApplying translations:")
    slide_index = 0
    
    for slide in prs.slides:
        slide_index += 1
        print(f"Processing Slide {slide_index}:")
        
        # Track shape index for this slide
        shape_index = 0
        
        for shape in slide.shapes:
            shape_index += 1
            shape_id = f"slide{slide_index}_shape{shape_index}"
            
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                print(f"  - Processing table in shape {shape_index}")
                
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            original_text = cell.text
                            
                            # Try exact text matching for table cells
                            if original_text in translations:
                                translation = translations[original_text]
                                print(f"    - Cell [{row_idx},{col_idx}] match: '{original_text}' -> '{translation}'")
                                cell.text = translation
                            else:
                                # Try fuzzy matching as fallback
                                for key, trans in translations.items():
                                    if (key.lower() in original_text.lower() or 
                                        original_text.lower() in key.lower() or
                                        original_text.strip() == key.strip()):
                                        translation = trans
                                        print(f"    - Cell [{row_idx},{col_idx}] fuzzy match: '{original_text}' -> '{translation}'")
                                        cell.text = translation
                                        break
            
            # Handle regular text
            elif hasattr(shape, "text") and shape.text:
                original_text = shape.text
                translation = None
                
                # Try exact text matching first (most reliable)
                if original_text in translations:
                    translation = translations[original_text]
                    print(f"  - Exact match: '{original_text}' -> '{translation}'")
                
                # Try matching paragraphs individually
                elif "\n" in original_text:
                    paragraphs = original_text.split("\n")
                    translated_parts = []
                    
                    for para in paragraphs:
                        if para.strip() in translations:
                            translated_parts.append(translations[para.strip()])
                            print(f"  - Paragraph match: '{para.strip()}' -> '{translations[para.strip()]}'")
                        else:
                            # Keep original if no translation
                            translated_parts.append(para)
                            print(f"  - No match for paragraph: '{para.strip()}'")
                    
                    if any(part != paragraphs[i] for i, part in enumerate(translated_parts)):
                        translation = "\n".join(translated_parts)
                
                # Last resort - try basic heuristic matching
                if not translation:
                    for key, trans in translations.items():
                        if key.lower() in original_text.lower() or original_text.lower() in key.lower():
                            translation = trans
                            print(f"  - Heuristic match: '{original_text}' to '{key}' -> '{translation}'")
                            break
                
                # Apply the translation if we found one
                if translation:
                    print(f"  - Applying translation for: '{original_text}'")
                    
                    # Replace text in the shape
                    if hasattr(shape, "text_frame"):
                        # Clear existing paragraphs
                        for i in range(len(shape.text_frame.paragraphs)-1, 0, -1):
                            p = shape.text_frame.paragraphs[i]
                            p.text = ""
                        
                        # Set the new text in the first paragraph
                        if len(shape.text_frame.paragraphs) > 0:
                            shape.text_frame.paragraphs[0].text = translation
                        
                        # Add additional paragraphs if needed
                        if "\n" in translation and len(translation.split("\n")) > 1:
                            paras = translation.split("\n")[1:]  # Skip first para as it's already set
                            for para_text in paras:
                                if para_text.strip():  # Only add non-empty paragraphs
                                    p = shape.text_frame.add_paragraph()
                                    p.text = para_text
                else:
                    print(f"  - No translation found for: '{original_text}'")
    
    # Save the result
    prs.save(sys.argv[3])
    print(f"\nTranslation successful! Output saved to {sys.argv[3]}")
    print(f"Translated to Japanese: 日本語に翻訳されました")
    
except Exception as e:
    print(f"Error during update: {str(e)}")
    sys.exit(1)