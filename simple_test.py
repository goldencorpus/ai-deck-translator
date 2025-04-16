#!/usr/bin/env python3
"""
Simple test script for slide notes.

This script creates a simple PowerPoint presentation with slide notes
and verifies that the notes are present.
"""
import os
from pptx import Presentation

# Create test directory if it doesn't exist
TEST_DIR = "test_output"
os.makedirs(TEST_DIR, exist_ok=True)

# File paths
TEST_PPTX = os.path.join(TEST_DIR, "simple_test.pptx")

def create_test_presentation():
    """Create a simple PowerPoint presentation with slide notes."""
    print("Creating test presentation...")
    prs = Presentation()
    
    # Slide 1: Title slide with notes
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Test Presentation"
    subtitle.text = "Testing Slide Notes"
    
    # Add notes to slide 1
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the title slide."
    
    # Save the presentation
    prs.save(TEST_PPTX)
    print(f"Test presentation created: {TEST_PPTX}")
    
    return prs

def verify_notes(prs):
    """Verify that the notes are present in the presentation."""
    print("\nVerifying notes...")
    
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            print(f"  Has notes: Yes")
            print(f"  Notes text: {notes_text}")
        else:
            print(f"  Has notes: No")

def main():
    """Main function."""
    # Create test presentation
    prs = create_test_presentation()
    
    # Verify notes
    verify_notes(prs)
    
    # Reopen the presentation and verify notes again
    print("\nReopening presentation...")
    prs = Presentation(TEST_PPTX)
    verify_notes(prs)

if __name__ == "__main__":
    main() 