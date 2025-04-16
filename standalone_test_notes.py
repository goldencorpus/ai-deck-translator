#!/usr/bin/env python3
"""
Standalone test script for slide notes translation feature.

This script creates a simple PowerPoint presentation with slide notes,
extracts text from the presentation, translates it, and updates the
presentation with the translated text.
"""
import os
import sys
from pptx import Presentation
import xml.etree.ElementTree as ET
import zipfile
import re
import shutil

# Create test directory if it doesn't exist
TEST_DIR = "test_output"
os.makedirs(TEST_DIR, exist_ok=True)

# File paths
TEST_PPTX = os.path.join(TEST_DIR, "test_presentation.pptx")
OUTPUT_PPTX = os.path.join(TEST_DIR, "test_presentation_translated.pptx")

# XML namespaces used in PPTX files
namespaces = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
}

def create_test_presentation():
    """Create a simple PowerPoint presentation with slide notes."""
    print("Creating test presentation...")
    prs = Presentation()
    
    # Slide 1: Title slide with notes
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Test Presentation"
    subtitle.text = "Testing Slide Notes Translation"
    
    # Add notes to slide 1
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the title slide. They should be translated."
    
    # Slide 2: Content slide with notes
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content Slide"
    content.text = "This is a bullet point\nThis is another bullet point\nThis is a third bullet point"
    
    # Add notes to slide 2
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the content slide. They contain important information for the presenter."
    
    # Slide 3: Content slide with notes
    slide_layout = prs.slide_layouts[1]  # Content slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Another Content Slide"
    content.text = "This slide has different content\nWith multiple bullet points\nAnd some more text"
    
    # Add notes to slide 3
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are more detailed speaker notes. They explain what the presenter should say during this slide."
    
    # Save the presentation
    prs.save(TEST_PPTX)
    print(f"Test presentation created: {TEST_PPTX}")

def extract_text(pptx_file):
    """Extract text from a PowerPoint presentation, including slide notes."""
    print(f"Extracting text from {pptx_file}...")
    
    # Dictionary to store text elements
    text_elements = {}
    
    # List to store slide metadata
    slide_metadata = []
    
    # Open the presentation
    prs = Presentation(pptx_file)
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1
        print(f"Processing slide {slide_number}...")
        
        # Create slide metadata
        metadata = {
            "slide_number": slide_number,
            "layout": slide.slide_layout.name if hasattr(slide.slide_layout, "name") else "Unknown",
            "elements": []
        }
        
        # Extract text from shapes
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = f"slide{slide_number}_shape{shape_idx}"
            
            # Extract text from text frames
            if hasattr(shape, "text") and shape.text.strip():
                text_elements[shape_id] = shape.text
                metadata["elements"].append({
                    "id": shape_id,
                    "type": "shape",
                    "text": shape.text
                })
            
            # Extract text from tables
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            cell_id = f"{shape_id}_table_r{row_idx}c{col_idx}"
                            text_elements[cell_id] = cell.text
                            metadata["elements"].append({
                                "id": cell_id,
                                "type": "table_cell",
                                "text": cell.text
                            })
        
        # Extract slide notes
        if slide.has_notes_slide and slide.notes_slide:
            notes_text = ""
            for shape in slide.notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    notes_text += shape.text + "\n"
            
            if notes_text.strip():
                metadata["notes"] = notes_text.strip()
                print(f"  Found notes: {notes_text.strip()}")
        
        # Add slide metadata to the list
        slide_metadata.append(metadata)
    
    return text_elements, slide_metadata

def mock_translate(text, target_language):
    """Mock translation function that adds a prefix to the text."""
    prefix_map = {
        "es": "[ES] ",
        "fr": "[FR] ",
        "de": "[DE] ",
        "ja": "[JA] ",
    }
    prefix = prefix_map.get(target_language, "[??] ")
    return prefix + text

def update_slides(pptx_file, output_file, translated_texts):
    """Update PowerPoint presentation with translated text."""
    print(f"Updating {pptx_file} with translated text...")
    
    # Create a copy of the presentation
    prs = Presentation(pptx_file)
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1
        print(f"Updating slide {slide_number}...")
        
        # Update text in shapes
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = f"slide{slide_number}_shape{shape_idx}"
            
            # Update text in text frames
            if hasattr(shape, "text") and shape.text.strip() and shape_id in translated_texts:
                # Replace text while preserving paragraph formatting
                if hasattr(shape, "text_frame") and shape.text != translated_texts[shape_id]:
                    # Clear existing paragraphs
                    existing_paragraphs = list(shape.text_frame.paragraphs)
                    
                    # If there's only one paragraph, simply update the text
                    if len(existing_paragraphs) == 1:
                        p = existing_paragraphs[0]
                        if p.runs:
                            p.runs[0].text = translated_texts[shape_id]
                        else:
                            p.text = translated_texts[shape_id]
                    else:
                        # For multiple paragraphs, try to match the structure
                        translated_lines = translated_texts[shape_id].split('\n')
                        
                        # Update existing paragraphs
                        for i, p in enumerate(existing_paragraphs):
                            if i < len(translated_lines):
                                if p.runs:
                                    p.runs[0].text = translated_lines[i]
                                else:
                                    p.text = translated_lines[i]
                            else:
                                # Clear extra paragraphs
                                if p.runs:
                                    p.runs[0].text = ""
                                else:
                                    p.text = ""
                        
                        # Add any additional paragraphs if needed
                        if len(translated_lines) > len(existing_paragraphs):
                            for i in range(len(existing_paragraphs), len(translated_lines)):
                                p = shape.text_frame.add_paragraph()
                                p.text = translated_lines[i]
            
            # Update text in tables
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_id = f"{shape_id}_table_r{row_idx}c{col_idx}"
                        if cell_id in translated_texts:
                            # Update cell text
                            cell.text = translated_texts[cell_id]
        
        # Update slide notes
        notes_id = f"slide{slide_number}_notes"
        if notes_id in translated_texts and slide.has_notes_slide and slide.notes_slide:
            # Find the notes text shape
            for notes_shape in slide.notes_slide.shapes:
                if hasattr(notes_shape, "text") and notes_shape.text.strip():
                    # Replace the notes text
                    if hasattr(notes_shape, "text_frame"):
                        # Clear existing paragraphs
                        while len(notes_shape.text_frame.paragraphs) > 0:
                            p = notes_shape.text_frame.paragraphs[0]
                            if hasattr(p, "element"):
                                p.element.getparent().remove(p.element)
                        
                        # Add the translated notes
                        p = notes_shape.text_frame.add_paragraph()
                        p.text = translated_texts[notes_id]
                        print(f"  Updated notes for slide {slide_number}")
                        break
    
    # Save the updated presentation
    prs.save(output_file)
    print(f"Saved updated presentation to {output_file}")
    
    # Now handle XML-level updates for notes
    update_xml_notes(pptx_file, output_file, translated_texts)
    
    return True

def update_xml_notes(original_file, updated_file, translated_texts):
    """Update slide notes in the PPTX file at the XML level."""
    print("Performing XML-level updates for notes...")
    
    # Create temporary directory for extraction
    temp_dir = os.path.join(os.path.dirname(updated_file), "_temp_pptx_update")
    os.makedirs(temp_dir, exist_ok=True)
    
    try:
        # Extract the updated file
        with zipfile.ZipFile(updated_file, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Process each slide
        slide_files = [f for f in os.listdir(os.path.join(temp_dir, "ppt", "slides")) if f.startswith("slide") and f.endswith(".xml")]
        
        for slide_file in slide_files:
            slide_number = int(re.search(r'slide(\d+)\.xml', slide_file).group(1))
            
            # Update slide notes in XML if needed
            notes_id = f"slide{slide_number}_notes"
            if notes_id in translated_texts:
                # Check if there's a notes slide XML file
                notes_slide_path = os.path.join(temp_dir, "ppt", "notesSlides", f"notesSlide{slide_number}.xml")
                if os.path.exists(notes_slide_path):
                    try:
                        # Parse the notes slide XML
                        notes_tree = ET.parse(notes_slide_path)
                        notes_root = notes_tree.getroot()
                        
                        # Find all text elements in the notes slide
                        text_elements = notes_root.findall('.//a:t', namespaces)
                        
                        # Split the translated notes into lines
                        translated_lines = translated_texts[notes_id].split('\n')
                        
                        # Update the text elements
                        if len(text_elements) == 1:
                            text_elements[0].text = translated_texts[notes_id]
                        elif len(text_elements) > 1:
                            # For multiple text elements, try to match the structure
                            for i, text_elem in enumerate(text_elements):
                                if i < len(translated_lines):
                                    text_elem.text = translated_lines[i]
                                else:
                                    text_elem.text = ""
                        
                        # Save the updated notes slide
                        notes_tree.write(notes_slide_path, encoding='utf-8', xml_declaration=True)
                        print(f"  Updated notes slide XML for slide {slide_number}")
                    except Exception as e:
                        print(f"  Error updating notes slide XML: {e}")
        
        # Recreate the PPTX file
        with zipfile.ZipFile(updated_file, 'w') as zip_ref:
            for root_dir, _, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zip_ref.write(file_path, arcname)
        
        print("XML-level updates completed successfully")
    except Exception as e:
        print(f"Error updating XML elements: {e}")
    finally:
        # Clean up temporary directory
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
            print("Cleaned up temporary directory")

def test_notes_translation():
    """Test the slide notes translation feature."""
    # Create test presentation
    create_test_presentation()
    
    # Extract text from the presentation
    text_elements, slide_metadata = extract_text(TEST_PPTX)
    
    # Print extracted text elements
    print("\nExtracted text elements:")
    for key, value in text_elements.items():
        print(f"{key}: {value}")
    
    # Print slide metadata (including notes)
    print("\nSlide metadata (including notes):")
    for slide in slide_metadata:
        print(f"Slide {slide['slide_number']}:")
        if 'notes' in slide and slide['notes']:
            print(f"  Notes: {slide['notes']}")
    
    # Translate text elements and notes
    print("\nTranslating text elements and notes...")
    target_language = "es"  # Spanish
    
    # Create dictionary of translated elements
    translated_elements = {}
    
    # Translate text elements
    for key, value in text_elements.items():
        translated_elements[key] = mock_translate(value, target_language)
    
    # Translate notes
    for slide in slide_metadata:
        if 'notes' in slide and slide['notes']:
            slide_number = slide['slide_number']
            notes_id = f"slide{slide_number}_notes"
            translated_elements[notes_id] = mock_translate(slide['notes'], target_language)
    
    # Print translated elements
    print("\nTranslated elements:")
    for key, value in translated_elements.items():
        print(f"{key}: {value}")
    
    # Update the presentation with translated text
    print("\nUpdating presentation with translated text...")
    success = update_slides(TEST_PPTX, OUTPUT_PPTX, translated_elements)
    
    if success:
        print(f"\nTranslation successful! Output file: {OUTPUT_PPTX}")
    else:
        print("\nTranslation failed!")

if __name__ == "__main__":
    test_notes_translation() 